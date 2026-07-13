from __future__ import annotations

import shutil
from pathlib import Path

SAMPLE_SEARCH_TASKS = ("预算审批", "项目里程碑", "负责人")


def unique_path(parent: Path, name: str) -> Path:
    candidate = parent / name
    number = 2
    while candidate.exists():
        candidate = parent / f"{name}-{number}"
        number += 1
    return candidate


def default_sample_paths() -> tuple[Path, Path]:
    documents = Path.home() / "Documents"
    raw = unique_path(documents, "Octopus 示例资料")
    index = unique_path(documents, "Octopus 示例索引")
    return raw, index


def _pdf_bytes() -> bytes:
    stream = b"BT /F1 18 Tf 72 720 Td (Octopus project meeting notes and budget approval) Tj ET"
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length "
        + str(len(stream)).encode("ascii")
        + b" >>\nstream\n"
        + stream
        + b"\nendstream",
    ]
    payload = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for number, value in enumerate(objects, start=1):
        offsets.append(len(payload))
        payload.extend(f"{number} 0 obj\n".encode("ascii"))
        payload.extend(value)
        payload.extend(b"\nendobj\n")
    xref = len(payload)
    payload.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    payload.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        payload.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    payload.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
        f"startxref\n{xref}\n%%EOF\n".encode("ascii")
    )
    return bytes(payload)


def materialize_sample_repository(destination: Path) -> Path:
    destination = destination.expanduser().resolve()
    if destination.exists():
        raise FileExistsError(f"Sample destination already exists: {destination}")
    destination.mkdir(parents=True, exist_ok=False)
    try:
        (destination / "项目说明.md").write_text(
            "# 星河项目\n\n"
            "预算审批截止日期是 2026 年 8 月 15 日。\n\n"
            "## 固定搜索任务\n\n"
            "1. 预算审批\n"
            "2. 项目里程碑\n"
            "3. 负责人\n",
            encoding="utf-8",
        )

        from docx import Document

        document = Document()
        document.add_heading("星河项目需求", level=1)
        document.add_paragraph("目标是建立本地、安全、可解释的资料索引。")
        document.add_paragraph("负责人：示例团队；预算审批：财务负责人。")
        document.save(str(destination / "项目需求.docx"))

        from openpyxl import Workbook  # type: ignore[import-untyped]

        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "预算"
        sheet.append(["项目", "预算", "审批状态"])
        sheet.append(["资料索引", 120000, "待审批"])
        workbook.save(destination / "项目预算.xlsx")

        from pptx import Presentation

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "星河项目里程碑"
        slide.placeholders[1].text = "需求确认\n预算审批\n首次交付"
        presentation.save(str(destination / "项目里程碑.pptx"))

        (destination / "会议纪要.pdf").write_bytes(_pdf_bytes())

        from PIL import Image, ImageDraw

        image = Image.new("RGB", (900, 180), "white")
        draw = ImageDraw.Draw(image)
        draw.text((40, 60), "OCTOPUS OCR SAMPLE - BUDGET APPROVAL 2026", fill="black")
        image.save(destination / "扫描记录.png")
    except Exception:
        shutil.rmtree(destination, ignore_errors=True)
        raise
    return destination
