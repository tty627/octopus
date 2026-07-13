from __future__ import annotations

import csv
import io
import zipfile
from collections import Counter
from pathlib import Path
from typing import Any

from .models import ContentParser, ExtractedDocument, ExtractionEvidence
from .utils import truncate

TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".py",
    ".js",
    ".ts",
    ".tsx",
    ".jsx",
    ".json",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".cfg",
    ".xml",
    ".html",
    ".css",
    ".sql",
    ".sh",
    ".ps1",
    ".bat",
    ".c",
    ".cpp",
    ".h",
    ".java",
    ".go",
    ".rs",
}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tif", ".tiff", ".webp"}
MAX_EXTRACTED_CHARACTERS = 80_000


def _evidence(
    locator: str,
    kind: str,
    text: str = "",
    *,
    method: str = "native",
    confidence: float | None = None,
) -> ExtractionEvidence:
    excerpt = " ".join(text.split())[:300]
    return ExtractionEvidence(
        locator=locator,
        kind=kind,
        text_excerpt=excerpt,
        extraction_method=method,
        confidence=confidence,
    )


def _document(
    *,
    name: str,
    document_type: str,
    text: str = "",
    structure: list[str] | None = None,
    metadata: dict[str, Any] | None = None,
    quality_flags: list[str] | None = None,
    evidence: list[ExtractionEvidence] | None = None,
    extraction_stats: dict[str, int | float | str | bool] | None = None,
    unsupported: bool = False,
) -> ExtractedDocument:
    return ExtractedDocument(
        name=name,
        document_type=document_type,
        text=truncate(text, MAX_EXTRACTED_CHARACTERS),
        structure=structure or [],
        metadata=metadata or {},
        quality_flags=sorted(set(quality_flags or [])),
        evidence=(evidence or [])[:200],
        extraction_stats=extraction_stats or {},
        unsupported=unsupported,
        text_characters=len(text),
        truncated=len(text) > MAX_EXTRACTED_CHARACTERS,
    )


def _package_member_count(path: Path, prefix: str) -> int:
    try:
        with zipfile.ZipFile(path) as archive:
            return sum(
                name.startswith(prefix) and not name.endswith("/") for name in archive.namelist()
            )
    except (OSError, zipfile.BadZipFile):
        return 0


def is_plain_text(path: Path) -> bool:
    if path.suffix.casefold() in TEXT_EXTENSIONS:
        return True
    if path.suffix.casefold() == ".csv":
        return path.stat().st_size <= 5 * 1024 * 1024
    return False


def _rapid_ocr_text(image: Any) -> tuple[str, list[str]]:
    try:
        from rapidocr import RapidOCR
    except ImportError:
        return "", ["ocr_unavailable"]
    try:
        result = RapidOCR()(image)
        texts: list[str] = []
        if hasattr(result, "txts"):
            texts = [str(item) for item in result.txts]
        elif isinstance(result, tuple) and result:
            rows = result[0] or []
            for row in rows:
                if isinstance(row, (list, tuple)) and len(row) >= 2:
                    value = row[1]
                    texts.append(str(value[0] if isinstance(value, (list, tuple)) else value))
        elif isinstance(result, list):
            for row in result:
                if isinstance(row, (list, tuple)) and len(row) >= 2:
                    value = row[1]
                    texts.append(str(value[0] if isinstance(value, (list, tuple)) else value))
        return "\n".join(texts), [] if texts else ["ocr_returned_no_text"]
    except Exception as error:  # OCR engines raise backend-specific exceptions.
        return "", [f"ocr_failed:{type(error).__name__}"]


class TextParser:
    def can_handle(self, path: Path) -> bool:
        return is_plain_text(path)

    def extract(self, path: Path) -> ExtractedDocument:
        text = path.read_text(encoding="utf-8-sig", errors="replace")
        structure = [line.strip() for line in text.splitlines() if line.lstrip().startswith("#")][
            :100
        ]
        if path.suffix.casefold() == ".csv":
            try:
                rows = list(csv.reader(io.StringIO(text[:50_000])))
                width = max((len(row) for row in rows), default=0)
                structure.insert(0, f"CSV sample: {len(rows)} rows, {width} columns")
            except csv.Error:
                pass
        nonempty = [line.strip() for line in text.splitlines() if line.strip()]
        evidence = [
            _evidence(f"heading:{index}", "heading", heading)
            for index, heading in enumerate(structure[:50], start=1)
        ]
        if not evidence and nonempty:
            evidence.append(_evidence("line:1", "text", nonempty[0]))
        return _document(
            name=path.name,
            document_type=path.suffix.casefold().lstrip(".") or "text",
            text=text,
            structure=structure,
            evidence=evidence,
            extraction_stats={"line_count": len(text.splitlines())},
        )


class PDFParser:
    def can_handle(self, path: Path) -> bool:
        return path.suffix.casefold() == ".pdf"

    def extract(self, path: Path) -> ExtractedDocument:
        try:
            import pypdfium2 as pdfium  # type: ignore[import-untyped]
            from pypdf import PdfReader
        except ImportError:
            return ExtractedDocument(
                name=path.name,
                document_type="pdf",
                quality_flags=["pdf_parser_unavailable"],
                unsupported=True,
            )
        texts: list[str] = []
        structure: list[str] = []
        flags: list[str] = []
        evidence: list[ExtractionEvidence] = []
        ocr_pages = 0
        ocr_attempts = 0
        native_text_pages = 0
        reader = PdfReader(str(path), strict=False)
        renderer = pdfium.PdfDocument(str(path))
        try:
            for page_number, page in enumerate(reader.pages, start=1):
                text = (page.extract_text() or "").strip()
                method = "native"
                if len(text) < 40 and ocr_attempts < 20:
                    ocr_attempts += 1
                    render_page = renderer[page_number - 1]
                    bitmap = render_page.render(scale=1.5)
                    image = bitmap.to_pil()
                    ocr_text, ocr_flags = _rapid_ocr_text(image)
                    bitmap.close()
                    render_page.close()
                    if ocr_text:
                        text = ocr_text
                        ocr_pages += 1
                        method = "ocr"
                    elif len(text) < 40:
                        flags.append(f"low_text_page:{page_number}")
                    flags.extend(ocr_flags)
                if method == "native" and text:
                    native_text_pages += 1
                texts.append(f"[Page {page_number}]\n{text}")
                first_line = next((line.strip() for line in text.splitlines() if line.strip()), "")
                structure.append(f"P{page_number}: {first_line[:120] or '[no text]'}")
                evidence.append(
                    _evidence(
                        f"page:{page_number}",
                        "page",
                        first_line,
                        method=method,
                    )
                )
        finally:
            renderer.close()
        metadata: dict[str, Any] = {
            str(key).lstrip("/"): str(value)
            for key, value in (reader.metadata or {}).items()
            if value
        }
        metadata["page_count"] = len(reader.pages)
        metadata["ocr_page_count"] = ocr_pages
        metadata["ocr_attempt_count"] = ocr_attempts
        return _document(
            name=path.name,
            document_type="pdf",
            text="\n\n".join(texts),
            structure=structure,
            metadata=metadata,
            quality_flags=flags,
            evidence=evidence,
            extraction_stats={
                "page_count": len(reader.pages),
                "native_text_pages": native_text_pages,
                "ocr_attempts": ocr_attempts,
                "ocr_pages": ocr_pages,
            },
        )


class DocxParser:
    def can_handle(self, path: Path) -> bool:
        return path.suffix.casefold() == ".docx"

    def extract(self, path: Path) -> ExtractedDocument:
        try:
            from docx import Document
        except ImportError:
            return ExtractedDocument(
                name=path.name,
                document_type="word",
                quality_flags=["docx_parser_unavailable"],
                unsupported=True,
            )
        document = Document(str(path))
        text: list[str] = []
        structure: list[str] = []
        evidence: list[ExtractionEvidence] = []
        for paragraph_index, paragraph in enumerate(document.paragraphs, start=1):
            value = paragraph.text.strip()
            if not value:
                continue
            text.append(value)
            if paragraph.style and paragraph.style.name.casefold().startswith("heading"):
                structure.append(f"{paragraph.style.name}: {value[:160]}")
                evidence.append(_evidence(f"paragraph:{paragraph_index}", "heading", value))
        for table_index, table in enumerate(document.tables, start=1):
            rows = []
            for row in table.rows[:20]:
                rows.append(" | ".join(cell.text.strip() for cell in row.cells[:20]))
            text.append(f"[Table {table_index}]\n" + "\n".join(rows))
            structure.append(
                f"Table {table_index}: {len(table.rows)} rows x {len(table.columns)} columns"
            )
            evidence.append(
                _evidence(
                    f"table:{table_index}",
                    "table",
                    rows[0] if rows else "",
                )
            )
        properties = document.core_properties
        metadata = {
            "paragraph_count": len(document.paragraphs),
            "table_count": len(document.tables),
            "inline_shape_count": len(document.inline_shapes),
            "media_file_count": _package_member_count(path, "word/media/"),
            "title": properties.title or "",
            "subject": properties.subject or "",
            "author": properties.author or "",
        }
        return _document(
            name=path.name,
            document_type="word",
            text="\n\n".join(text),
            structure=structure,
            metadata=metadata,
            evidence=evidence,
            extraction_stats={
                "paragraphs": len(document.paragraphs),
                "tables": len(document.tables),
                "headings": sum(item.kind == "heading" for item in evidence),
            },
        )


class XlsxParser:
    def can_handle(self, path: Path) -> bool:
        return path.suffix.casefold() in {".xlsx", ".xlsm"}

    def extract(self, path: Path) -> ExtractedDocument:
        try:
            import openpyxl  # type: ignore[import-untyped]
        except ImportError:
            return ExtractedDocument(
                name=path.name,
                document_type="excel",
                quality_flags=["xlsx_parser_unavailable"],
                unsupported=True,
            )
        workbook = openpyxl.load_workbook(path, read_only=True, data_only=False)
        structure: list[str] = []
        samples: list[str] = []
        evidence: list[ExtractionEvidence] = []
        formula_count = 0
        sampled_cells = 0
        hidden_sheets = 0
        chart_count = _package_member_count(path, "xl/charts/")
        for sheet in workbook.worksheets:
            if sheet.sheet_state != "visible":
                hidden_sheets += 1
            structure.append(
                f"Sheet {sheet.title}: rows={sheet.max_row}, "
                f"columns={sheet.max_column}, state={sheet.sheet_state}"
            )
            sample_rows = []
            for row in sheet.iter_rows(
                min_row=1,
                max_row=min(sheet.max_row, 50),
                max_col=min(sheet.max_column, 30),
                values_only=True,
            ):
                sampled_cells += len(row)
                formula_count += sum(
                    isinstance(value, str) and value.startswith("=") for value in row
                )
                sample_rows.append(
                    " | ".join("" if value is None else str(value) for value in row[:20])
                )
            samples.append(f"[Sheet: {sheet.title}]\n" + "\n".join(sample_rows))
            evidence.append(
                _evidence(
                    f"sheet:{sheet.title}",
                    "worksheet",
                    sample_rows[0] if sample_rows else "",
                )
            )
        sheet_names = list(workbook.sheetnames)
        workbook.close()
        return _document(
            name=path.name,
            document_type="excel",
            text="\n\n".join(samples),
            structure=structure,
            metadata={
                "sheet_names": sheet_names,
                "hidden_sheet_count": hidden_sheets,
                "sampled_formula_count": formula_count,
                "chart_count": chart_count,
            },
            evidence=evidence,
            extraction_stats={
                "sheets": len(sheet_names),
                "sampled_cells": sampled_cells,
                "sampled_formulas": formula_count,
                "hidden_sheets": hidden_sheets,
                "charts": chart_count,
            },
        )


class PptxParser:
    def can_handle(self, path: Path) -> bool:
        return path.suffix.casefold() == ".pptx"

    def extract(self, path: Path) -> ExtractedDocument:
        try:
            from pptx import Presentation
        except ImportError:
            return ExtractedDocument(
                name=path.name,
                document_type="ppt",
                quality_flags=["pptx_parser_unavailable"],
                unsupported=True,
            )
        presentation = Presentation(str(path))
        texts: list[str] = []
        structure: list[str] = []
        evidence: list[ExtractionEvidence] = []
        shape_count = 0
        notes_count = 0
        shape_types: Counter[str] = Counter()
        for number, slide in enumerate(presentation.slides, start=1):
            shape_count += len(slide.shapes)
            shape_types.update(str(shape.shape_type) for shape in slide.shapes)
            slide_text = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            title = slide_text[0] if slide_text else "[no title]"
            structure.append(f"Slide {number}: {title[:160]}")
            texts.append(f"[Slide {number}]\n" + "\n".join(slide_text))
            evidence.append(_evidence(f"slide:{number}", "slide", title))
            try:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    notes_count += 1
                    texts.append(f"[Slide {number} notes]\n{notes}")
                    evidence.append(_evidence(f"slide:{number}:notes", "speaker_notes", notes))
            except (AttributeError, KeyError):
                pass
        return _document(
            name=path.name,
            document_type="ppt",
            text="\n\n".join(texts),
            structure=structure,
            metadata={
                "slide_count": len(presentation.slides),
                "shape_count": shape_count,
                "slides_with_notes": notes_count,
                "media_file_count": _package_member_count(path, "ppt/media/"),
                "shape_types": dict(shape_types),
            },
            evidence=evidence,
            extraction_stats={
                "slides": len(presentation.slides),
                "shapes": shape_count,
                "slides_with_notes": notes_count,
            },
        )


class ImageParser:
    def can_handle(self, path: Path) -> bool:
        return path.suffix.casefold() in IMAGE_EXTENSIONS

    def extract(self, path: Path) -> ExtractedDocument:
        try:
            from PIL import Image
        except ImportError:
            return ExtractedDocument(
                name=path.name,
                document_type="image",
                quality_flags=["image_parser_unavailable"],
                unsupported=True,
            )
        with Image.open(path) as image:
            exif = image.getexif()
            frame_count = int(getattr(image, "n_frames", 1))
            metadata = {
                "width": image.width,
                "height": image.height,
                "mode": image.mode,
                "format": image.format,
                "frame_count": frame_count,
                "exif_tag_count": len(exif),
            }
            image.load()
            text, flags = _rapid_ocr_text(image)
        evidence = (
            [
                _evidence(
                    "image:ocr",
                    "ocr_text",
                    text,
                    method="ocr",
                )
            ]
            if text
            else []
        )
        return _document(
            name=path.name,
            document_type="image",
            text=text,
            structure=[f"Image: {metadata['width']}x{metadata['height']} {metadata['mode']}"],
            metadata=metadata,
            quality_flags=flags,
            evidence=evidence,
            extraction_stats={
                "ocr_characters": len(text),
                "ocr_succeeded": bool(text),
                "frames": frame_count,
            },
        )


class UnsupportedParser:
    def can_handle(self, path: Path) -> bool:
        return True

    def extract(self, path: Path) -> ExtractedDocument:
        return _document(
            name=path.name,
            document_type=path.suffix.casefold().lstrip(".") or "unknown",
            metadata={"size_bytes": path.stat().st_size},
            quality_flags=["unsupported_content_parser"],
            unsupported=True,
        )


class ParserRegistry:
    def __init__(self, parsers: list[ContentParser] | None = None) -> None:
        self.parsers = parsers or [
            PDFParser(),
            DocxParser(),
            XlsxParser(),
            PptxParser(),
            ImageParser(),
            TextParser(),
            UnsupportedParser(),
        ]

    def parser_for(self, path: Path) -> ContentParser:
        return next(parser for parser in self.parsers if parser.can_handle(path))

    def extract(self, path: Path) -> ExtractedDocument:
        return self.parser_for(path).extract(path)
