from __future__ import annotations

import codecs
import hashlib
import json
import math
import os
import re
import sqlite3
import time
import unicodedata
import uuid
from collections import Counter
from collections.abc import Callable, Sequence
from contextlib import closing
from contextvars import ContextVar
from dataclasses import dataclass, field, replace
from datetime import UTC, datetime
from difflib import SequenceMatcher
from functools import lru_cache
from pathlib import Path
from typing import Any, Literal, cast

from pydantic import Field

from .config import (
    global_config_lock,
    load_global_config,
    load_repository_config,
    save_global_config,
    workspace_storage_path,
    workspace_tasks_path,
)
from .credentials import CredentialStoreError, resolve_ai_api_key
from .models import AIConfig, GlobalWorkspace, OctopusModel, utc_now
from .utils import load_json, sha256_file
from .workspace_sources import (
    ArchiveCandidate,
    ArchivePolicy,
    EvidenceLocator,
    SourceRef,
    cache_expiry,
    materialize_source_ref,
    physical_source_ref,
    scan_archive,
)

WORKSPACE_SCHEMA_VERSION = "2.1"
PARSER_API_VERSION = "2.1"
READABLE_THRESHOLD = 0.72
PARTIAL_THRESHOLD = 0.45
MAX_TEXT_BYTES = 100 * 1024 * 1024
MAX_TEXT_CHARACTERS = 2_000_000
MAX_DOCX_PARAGRAPHS = 50_000
MAX_DOCX_TABLE_CELLS = 100_000
MAX_XLSX_CELLS = 100_000
MAX_PPTX_SLIDES = 5_000
MAX_IMAGE_PIXELS = 80_000_000
TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".rst",
    ".csv",
    ".json",
    ".jsonl",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".cfg",
    ".log",
    ".py",
    ".js",
    ".ts",
    ".tsx",
    ".jsx",
    ".html",
    ".css",
    ".sql",
}
OFFICE_EXTENSIONS = {".docx", ".xlsx", ".xlsm", ".pptx"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".webp", ".bmp"}
IGNORED_DIRECTORY_NAMES = {
    ".git",
    ".hg",
    ".svn",
    ".octopus",
    "node_modules",
    "__pycache__",
}

ExtractionProgressCallback = Callable[[dict[str, Any]], None]
_EXTRACTION_PROGRESS_CALLBACK: ContextVar[ExtractionProgressCallback | None] = ContextVar(
    "octopus_extraction_progress_callback",
    default=None,
)


class WorkspaceEvidence(OctopusModel):
    page_number: int | None = None
    locator: EvidenceLocator | None = None
    heading: str = ""
    excerpt: str
    reason: str
    quality_score: float = Field(ge=0.0, le=1.0)


class WorkspaceSearchResult(OctopusModel):
    document_id: str
    name: str
    relative_path: str
    extension: str
    content_hash: str
    size_bytes: int = Field(ge=0)
    modified_at: str
    page_count: int = Field(ge=0)
    readability: Literal["readable", "partial", "low"]
    readability_score: float = Field(ge=0.0, le=1.0)
    indexing_state: Literal["indexed", "metadata_only", "failed"]
    source_uri: str
    source_ref: SourceRef | None = None
    locator: EvidenceLocator | None = None
    quality_flags: list[str] = Field(default_factory=list)
    error_code: str = ""
    parser_key: str = ""
    parser_version: str = ""
    freshness_status: Literal[
        "current", "stale", "changed", "missing", "unverified", "unavailable", "needs_review"
    ] = "current"
    overview: str = ""
    best_evidence: WorkspaceEvidence
    additional_evidence: list[WorkspaceEvidence] = Field(default_factory=list)
    rank: int = Field(ge=1)


class WorkspaceSearchReport(OctopusModel):
    query: str
    requested_mode: Literal["local", "assisted"] = "local"
    actual_mode: Literal["local", "assisted", "degraded"] = "local"
    degradation_reason: str = ""
    answer: str = ""
    results: list[WorkspaceSearchResult] = Field(default_factory=list)
    candidate_count: int = Field(default=0, ge=0)
    duration_ms: int = Field(default=0, ge=0)


class AssistedSearchOrder(OctopusModel):
    ordered_document_ids: list[str] = Field(default_factory=list)
    answer: str = ""


class WorkspaceHealth(OctopusModel):
    document_count: int = 0
    readable_count: int = 0
    partial_count: int = 0
    low_quality_count: int = 0
    metadata_only_count: int = 0
    failed_count: int = 0
    last_sync_at: str = ""


class WorkspacePayload(OctopusModel):
    workspace_id: str
    name: str
    raw_path: str
    available: bool
    enabled: bool
    vision_enabled: bool
    legacy_index_present: bool
    health: WorkspaceHealth = Field(default_factory=WorkspaceHealth)


class WorkspaceDocument(OctopusModel):
    document_id: str
    name: str
    relative_path: str
    extension: str
    content_hash: str
    size_bytes: int = Field(ge=0)
    modified_at: str
    title: str
    overview: str = ""
    page_count: int = Field(ge=0)
    readability: Literal["readable", "partial", "low"]
    readability_score: float = Field(ge=0.0, le=1.0)
    indexing_state: Literal["indexed", "metadata_only", "failed"]
    error: str = ""
    source_uri: str
    source_ref: SourceRef | None = None
    locator: EvidenceLocator | None = None
    quality_flags: list[str] = Field(default_factory=list)
    error_code: str = ""
    parser_key: str = ""
    parser_version: str = ""
    freshness_status: Literal[
        "current", "stale", "changed", "missing", "unverified", "unavailable", "needs_review"
    ] = "current"


@dataclass(frozen=True)
class ExtractedPage:
    page_number: int | None
    text: str
    extraction_method: str
    quality_score: float
    locator: EvidenceLocator | None = None

    @property
    def readability(self) -> Literal["readable", "partial", "low"]:
        return readability_label(self.quality_score)


@dataclass(frozen=True)
class ExtractedSource:
    title: str
    pages: list[ExtractedPage]
    page_count: int
    status: str = "indexed"
    error: str = ""
    quality_flags: list[str] = field(default_factory=list)
    error_code: str = ""
    parser_key: str = ""
    parser_version: str = PARSER_API_VERSION


def _script_name(character: str) -> str:
    if character.isascii():
        return "LATIN" if character.isalpha() else "COMMON"
    name = unicodedata.name(character, "")
    for script in (
        "CJK",
        "HIRAGANA",
        "KATAKANA",
        "HANGUL",
        "GREEK",
        "CYRILLIC",
        "ARABIC",
        "HEBREW",
        "DEVANAGARI",
        "THAI",
        "BENGALI",
        "TIBETAN",
        "SINHALA",
        "ETHIOPIC",
        "GEORGIAN",
        "ARMENIAN",
        "LAO",
        "MYANMAR",
        "KHMER",
    ):
        if script in name:
            return script
    return "COMMON" if not character.isalpha() else "OTHER"


def readability_score(text: str) -> float:
    compact = [character for character in text if not character.isspace()]
    if not compact:
        return 0.0
    total = len(compact)
    printable_ratio = sum(character.isprintable() for character in compact) / total
    control_ratio = (
        sum(unicodedata.category(character).startswith("C") for character in compact) / total
    )
    replacement_ratio = sum(character in {"\ufffd", "\x00"} for character in compact) / total
    semantic_ratio = (
        sum(character.isalnum() or "\u4e00" <= character <= "\u9fff" for character in compact)
        / total
    )
    scripts = Counter(
        _script_name(character)
        for character in compact
        if character.isalpha() or "\u4e00" <= character <= "\u9fff"
    )
    meaningful_scripts = [
        name for name, count in scripts.items() if name != "COMMON" and count >= 2
    ]
    script_penalty = max(0, len(meaningful_scripts) - 3) * 0.14
    dominant_ratio = max(scripts.values(), default=total) / max(1, sum(scripts.values()))
    fragmented_script_penalty = 0.0
    if len(meaningful_scripts) >= 4 and dominant_ratio < 0.72:
        fragmented_script_penalty = 0.22
    repeated_ratio = max(Counter(compact).values()) / total
    repeated_penalty = max(0.0, repeated_ratio - 0.32) * 0.8
    repeated_ngram_ratio = 0.0
    if total >= 8 and len(set(compact)) / total < 0.45:
        repeated_ngram_ratio = max(
            (
                max(
                    Counter(
                        tuple(compact[index : index + size]) for index in range(total - size + 1)
                    ).values()
                )
                * size
                / total
            )
            for size in (2, 3, 4)
        )
    repeated_ngram_penalty = max(0.0, repeated_ngram_ratio - 0.40) * 0.75
    short_penalty = 0.12 if total < 20 else 0.0
    score = (
        0.42 * printable_ratio
        + 0.38 * min(1.0, semantic_ratio / 0.55)
        + 0.20 * min(1.0, total / 120)
        - control_ratio * 4.0
        - replacement_ratio * 5.0
        - script_penalty
        - fragmented_script_penalty
        - repeated_penalty
        - repeated_ngram_penalty
        - short_penalty
    )
    return round(max(0.0, min(1.0, score)), 4)


def readability_label(score: float) -> Literal["readable", "partial", "low"]:
    if score >= READABLE_THRESHOLD:
        return "readable"
    if score >= PARTIAL_THRESHOLD:
        return "partial"
    return "low"


def _readability_value(value: object) -> Literal["readable", "partial", "low"]:
    text = str(value)
    if text not in {"readable", "partial", "low"}:
        return "low"
    return cast(Literal["readable", "partial", "low"], text)


def _indexing_state_value(value: object) -> Literal["indexed", "metadata_only", "failed"]:
    text = str(value)
    if text not in {"indexed", "metadata_only", "failed"}:
        return "failed"
    return cast(Literal["indexed", "metadata_only", "failed"], text)


def normalize_search_text(value: str) -> str:
    normalized = unicodedata.normalize("NFKC", value).casefold()
    return " ".join(re.findall(r"[\w\u4e00-\u9fff]+", normalized, flags=re.UNICODE))


def search_terms(value: str) -> list[str]:
    normalized = normalize_search_text(value)
    terms = set(normalized.split())
    for block in re.findall(r"[\u4e00-\u9fff]+", normalized):
        for size in (2, 3):
            terms.update(block[index : index + size] for index in range(len(block) - size + 1))
    return sorted(term for term in terms if term)


def _fts_text(value: str) -> str:
    normalized = normalize_search_text(value)
    tokens = search_terms(value)
    return " ".join([normalized, *tokens]).strip()


def _excerpt(text: str, query: str, limit: int = 240) -> str:
    compact = " ".join(text.split())
    if not compact:
        return "正文识别质量较低，可按文件名查找。"
    normalized = compact.casefold()
    position = normalized.find(query.casefold())
    if position < 0:
        for term in search_terms(query):
            position = normalized.find(term)
            if position >= 0:
                break
    start = max(0, position - 70) if position >= 0 else 0
    value = compact[start : start + limit]
    if start:
        value = "…" + value
    if start + limit < len(compact):
        value += "…"
    return value


@lru_cache(maxsize=1)
def _ocr_engine() -> Any:
    from rapidocr import RapidOCR

    return RapidOCR()


def _ocr_text(image: Any) -> str:
    try:
        result = _ocr_engine()(image)
    except Exception:
        return ""
    if hasattr(result, "txts"):
        return "\n".join(str(item) for item in (result.txts or []))
    if isinstance(result, tuple) and result:
        rows = result[0] or []
    elif isinstance(result, list):
        rows = result
    else:
        rows = []
    texts: list[str] = []
    for row in rows:
        if isinstance(row, (list, tuple)) and len(row) >= 2:
            value = row[1]
            texts.append(str(value[0] if isinstance(value, (list, tuple)) else value))
    return "\n".join(texts)


def _extract_pdf(
    path: Path,
    progress_callback: ExtractionProgressCallback | None = None,
) -> ExtractedSource:
    import pypdfium2 as pdfium  # type: ignore[import-untyped]
    from pypdf import PdfReader

    document = pdfium.PdfDocument(str(path))
    reader: PdfReader | None = None
    pages: list[ExtractedPage] = []
    page_count = len(document)
    ocr_pages_completed = 0

    def report_progress(**changes: Any) -> None:
        if progress_callback is not None:
            progress_callback(dict(changes))

    try:
        for page_index in range(page_count):
            page_number = page_index + 1
            base_progress = {
                "current_page": page_number,
                "page_count": page_count,
                "pages_completed": page_index,
                "ocr_pages_completed": ocr_pages_completed,
            }
            report_progress(extraction_stage="pdfium", **base_progress)
            page = document[page_index]
            text_page = page.get_textpage()
            try:
                pdfium_text = text_page.get_text_range().strip()
            finally:
                text_page.close()
            selected_text = pdfium_text
            method = "pdfium"
            score = readability_score(selected_text)
            if score < READABLE_THRESHOLD:
                report_progress(extraction_stage="pypdf", **base_progress)
                try:
                    reader = reader or PdfReader(str(path), strict=False)
                    pypdf_text = (
                        (reader.pages[page_index].extract_text() or "").strip()
                        if page_index < len(reader.pages)
                        else ""
                    )
                except Exception:
                    pypdf_text = ""
                pypdf_score = readability_score(pypdf_text)
                if pypdf_score > score:
                    selected_text = pypdf_text
                    score = pypdf_score
                    method = "pypdf"
            if score < READABLE_THRESHOLD:
                report_progress(extraction_stage="ocr", **base_progress)
                bitmap = page.render(scale=2.0)
                try:
                    ocr_text = _ocr_text(bitmap.to_pil()).strip()
                finally:
                    bitmap.close()
                ocr_pages_completed += 1
                ocr_score = readability_score(ocr_text)
                if ocr_text and (ocr_score >= score + 0.05 or score < PARTIAL_THRESHOLD):
                    selected_text = ocr_text
                    score = ocr_score
                    method = "ocr"
            pages.append(
                ExtractedPage(
                    page_number=page_number,
                    text=selected_text,
                    extraction_method=method,
                    quality_score=score,
                    locator=EvidenceLocator(kind="page", page_number=page_number),
                )
            )
            report_progress(
                extraction_stage="page_complete",
                current_page=page_number,
                page_count=page_count,
                pages_completed=page_number,
                ocr_pages_completed=ocr_pages_completed,
            )
            page.close()
    finally:
        document.close()
    metadata: Any = reader.metadata if reader is not None else None
    metadata = metadata or {}
    title = str(metadata.get("/Title", "")).strip() or path.stem
    return ExtractedSource(
        title=title,
        pages=pages,
        page_count=len(pages),
        parser_key="v2.pdf",
        parser_version="2",
    )


def _text_encoding(sample: bytes) -> str:
    if sample.startswith(codecs.BOM_UTF8):
        return "utf-8-sig"
    if sample.startswith((codecs.BOM_UTF16_LE, codecs.BOM_UTF16_BE)):
        return "utf-16"
    for encoding in ("utf-8", "utf-16-le", "utf-16-be"):
        try:
            sample.decode(encoding)
            return encoding
        except UnicodeDecodeError:
            continue
    try:
        from charset_normalizer import from_bytes

        best = from_bytes(sample).best()
        if best is not None and best.encoding:
            return str(best.encoding)
    except ImportError:
        pass
    return "utf-8"


def _decode_text(path: Path) -> tuple[str, bool]:
    with path.open("rb") as source:
        sample = source.read(min(1024 * 1024, MAX_TEXT_BYTES))
        encoding = _text_encoding(sample)
        source.seek(0)
        decoder = codecs.getincrementaldecoder(encoding)(errors="replace")
        chunks: list[str] = []
        bytes_read = 0
        characters = 0
        truncated = False
        while bytes_read < MAX_TEXT_BYTES and characters < MAX_TEXT_CHARACTERS:
            chunk = source.read(min(1024 * 1024, MAX_TEXT_BYTES - bytes_read))
            if not chunk:
                break
            bytes_read += len(chunk)
            decoded = decoder.decode(chunk)
            remaining = MAX_TEXT_CHARACTERS - characters
            if len(decoded) > remaining:
                decoded = decoded[:remaining]
                truncated = True
            chunks.append(decoded)
            characters += len(decoded)
        if source.read(1):
            truncated = True
        if characters < MAX_TEXT_CHARACTERS:
            tail = decoder.decode(b"", final=True)
            remaining = MAX_TEXT_CHARACTERS - characters
            if len(tail) > remaining:
                tail = tail[:remaining]
                truncated = True
            chunks.append(tail)
    return "".join(chunks), truncated


def _extract_text(path: Path) -> ExtractedSource:
    text, truncated = _decode_text(path)
    lines = text.count("\n") + (1 if text else 0)
    return ExtractedSource(
        title=path.stem,
        pages=[
            ExtractedPage(
                page_number=None,
                text=text,
                extraction_method="text",
                quality_score=readability_score(text),
                locator=EvidenceLocator(
                    kind="text_line",
                    line_start=1 if lines else None,
                    line_end=lines or None,
                ),
            )
        ],
        page_count=0,
        quality_flags=["extraction_truncated"] if truncated else [],
        error_code="extraction_budget_exceeded" if truncated else "",
        parser_key="v2.text",
        parser_version="2",
    )


def _extract_docx(path: Path) -> ExtractedSource:
    try:
        from docx import Document
    except ImportError:
        return ExtractedSource(
            title=path.stem,
            pages=[],
            page_count=0,
            status="metadata_only",
            quality_flags=["docx_parser_unavailable"],
            error_code="docx_parser_unavailable",
            parser_key="v2.docx",
            parser_version="1",
        )
    document = Document(str(path))
    pages: list[ExtractedPage] = []
    flags: list[str] = []
    for index, paragraph in enumerate(document.paragraphs, start=1):
        if index > MAX_DOCX_PARAGRAPHS:
            flags.append("docx_paragraph_limit")
            break
        text = paragraph.text.strip()
        if not text:
            continue
        label = ""
        if paragraph.style and paragraph.style.name.casefold().startswith("heading"):
            label = paragraph.style.name
        pages.append(
            ExtractedPage(
                page_number=None,
                text=text,
                extraction_method="docx",
                quality_score=readability_score(text),
                locator=EvidenceLocator(
                    kind="paragraph",
                    paragraph_index=index,
                    label=label,
                ),
            )
        )
    cells = 0
    for table_index, table in enumerate(document.tables, start=1):
        rows: list[str] = []
        for row in table.rows:
            values: list[str] = []
            for cell in row.cells:
                cells += 1
                if cells > MAX_DOCX_TABLE_CELLS:
                    flags.append("docx_table_cell_limit")
                    break
                values.append(cell.text.strip())
            if values:
                rows.append(" | ".join(values))
            if cells > MAX_DOCX_TABLE_CELLS:
                break
        text = "\n".join(rows).strip()
        if text:
            pages.append(
                ExtractedPage(
                    page_number=None,
                    text=text,
                    extraction_method="docx-table",
                    quality_score=readability_score(text),
                    locator=EvidenceLocator(kind="table", table_index=table_index),
                )
            )
        if cells > MAX_DOCX_TABLE_CELLS:
            break
    properties = document.core_properties
    title = str(properties.title or "").strip() or path.stem
    return ExtractedSource(
        title=title,
        pages=pages,
        page_count=0,
        quality_flags=sorted(set(flags)),
        error_code="extraction_budget_exceeded" if flags else "",
        parser_key="v2.docx",
        parser_version="1",
    )


def _extract_xlsx(path: Path) -> ExtractedSource:
    try:
        import openpyxl  # type: ignore[import-untyped]
        from openpyxl.utils import get_column_letter  # type: ignore[import-untyped]
    except ImportError:
        return ExtractedSource(
            title=path.stem,
            pages=[],
            page_count=0,
            status="metadata_only",
            quality_flags=["xlsx_parser_unavailable"],
            error_code="xlsx_parser_unavailable",
            parser_key="v2.xlsx",
            parser_version="1",
        )
    workbook = openpyxl.load_workbook(
        path,
        read_only=True,
        data_only=False,
        keep_vba=path.suffix.casefold() == ".xlsm",
    )
    pages: list[ExtractedPage] = []
    flags: list[str] = []
    cells = 0
    try:
        for sheet in workbook.worksheets:
            rows: list[str] = []
            last_row = 0
            last_column = 0
            for row_number, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                values: list[str] = []
                for column_number, value in enumerate(row, start=1):
                    cells += 1
                    if cells > MAX_XLSX_CELLS:
                        flags.append("xlsx_cell_limit")
                        break
                    if value is not None:
                        last_row = row_number
                        last_column = max(last_column, column_number)
                    values.append("" if value is None else str(value))
                if any(values):
                    rows.append(" | ".join(values).rstrip())
                if cells > MAX_XLSX_CELLS:
                    break
            text = "\n".join(rows).strip()
            if text:
                cell_range = (
                    f"A1:{get_column_letter(last_column)}{last_row}"
                    if last_row and last_column
                    else ""
                )
                pages.append(
                    ExtractedPage(
                        page_number=None,
                        text=text,
                        extraction_method="xlsx",
                        quality_score=readability_score(text),
                        locator=EvidenceLocator(
                            kind="sheet",
                            sheet_name=str(sheet.title),
                            cell_range=cell_range,
                        ),
                    )
                )
            if cells > MAX_XLSX_CELLS:
                break
    finally:
        workbook.close()
    return ExtractedSource(
        title=path.stem,
        pages=pages,
        page_count=0,
        quality_flags=sorted(set(flags)),
        error_code="extraction_budget_exceeded" if flags else "",
        parser_key="v2.xlsx",
        parser_version="1",
    )


def _extract_pptx(path: Path) -> ExtractedSource:
    try:
        from pptx import Presentation
    except ImportError:
        return ExtractedSource(
            title=path.stem,
            pages=[],
            page_count=0,
            status="metadata_only",
            quality_flags=["pptx_parser_unavailable"],
            error_code="pptx_parser_unavailable",
            parser_key="v2.pptx",
            parser_version="1",
        )
    presentation = Presentation(str(path))
    pages: list[ExtractedPage] = []
    flags: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        if slide_number > MAX_PPTX_SLIDES:
            flags.append("pptx_slide_limit")
            break
        values = [
            str(shape.text).strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and str(shape.text).strip()
        ]
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except (AttributeError, KeyError):
            notes = ""
        if notes:
            values.append(f"[Notes]\n{notes}")
        text = "\n".join(values).strip()
        if not text:
            continue
        pages.append(
            ExtractedPage(
                page_number=None,
                text=text,
                extraction_method="pptx",
                quality_score=readability_score(text),
                locator=EvidenceLocator(kind="slide", slide_number=slide_number),
            )
        )
    return ExtractedSource(
        title=path.stem,
        pages=pages,
        page_count=0,
        quality_flags=sorted(set(flags)),
        error_code="extraction_budget_exceeded" if flags else "",
        parser_key="v2.pptx",
        parser_version="1",
    )


def _extract_image(path: Path) -> ExtractedSource:
    try:
        from PIL import Image
    except ImportError:
        return ExtractedSource(
            title=path.stem,
            pages=[],
            page_count=0,
            status="metadata_only",
            quality_flags=["image_parser_unavailable"],
            error_code="image_parser_unavailable",
            parser_key="v2.image",
            parser_version="1",
        )
    with Image.open(path) as image:
        if image.width * image.height > MAX_IMAGE_PIXELS:
            return ExtractedSource(
                title=path.stem,
                pages=[],
                page_count=0,
                status="metadata_only",
                quality_flags=["image_pixel_limit"],
                error_code="image_pixel_limit",
                parser_key="v2.image",
                parser_version="1",
            )
        image.load()
        text = _ocr_text(image).strip()
    flags = [] if text else ["ocr_returned_no_text"]
    return ExtractedSource(
        title=path.stem,
        pages=[
            ExtractedPage(
                page_number=None,
                text=text,
                extraction_method="ocr",
                quality_score=readability_score(text),
                locator=EvidenceLocator(kind="image", label=path.name),
            )
        ],
        page_count=0,
        quality_flags=flags,
        error_code="ocr_returned_no_text" if flags else "",
        parser_key="v2.image",
        parser_version="1",
    )


@dataclass(frozen=True)
class WorkspaceParser:
    key: str
    version: str
    extensions: frozenset[str]
    extract: Callable[[Path], ExtractedSource]


class WorkspaceParserRegistry:
    def __init__(self) -> None:
        self.parsers = [
            WorkspaceParser(
                "v2.pdf",
                "2",
                frozenset({".pdf"}),
                lambda path: _extract_pdf(path, _EXTRACTION_PROGRESS_CALLBACK.get()),
            ),
            WorkspaceParser("v2.docx", "1", frozenset({".docx"}), _extract_docx),
            WorkspaceParser("v2.xlsx", "1", frozenset({".xlsx", ".xlsm"}), _extract_xlsx),
            WorkspaceParser("v2.pptx", "1", frozenset({".pptx"}), _extract_pptx),
            WorkspaceParser("v2.image", "1", frozenset(IMAGE_EXTENSIONS), _extract_image),
            WorkspaceParser("v2.text", "2", frozenset(TEXT_EXTENSIONS), _extract_text),
        ]

    def parser_for(self, path: Path) -> WorkspaceParser | None:
        suffix = path.suffix.casefold()
        return next((parser for parser in self.parsers if suffix in parser.extensions), None)

    def signature_for(self, path: Path) -> tuple[str, str]:
        parser = self.parser_for(path)
        return (parser.key, parser.version) if parser else ("v2.unsupported", "1")

    def extract_source(self, path: Path) -> ExtractedSource:
        parser = self.parser_for(path)
        if parser is None:
            return ExtractedSource(
                title=path.stem,
                pages=[],
                page_count=0,
                status="metadata_only",
                quality_flags=["unsupported_content_parser"],
                error_code="unsupported_content_parser",
                parser_key="v2.unsupported",
                parser_version="1",
            )
        extracted = parser.extract(path)
        if extracted.parser_key == parser.key and extracted.parser_version == parser.version:
            return extracted
        return replace(extracted, parser_key=parser.key, parser_version=parser.version)


PARSER_REGISTRY = WorkspaceParserRegistry()


def parser_signature(path: Path) -> tuple[str, str]:
    if path.suffix.casefold() == ".zip":
        return ("v2.archive", "1")
    return PARSER_REGISTRY.signature_for(path)


def extract_source(path: Path) -> ExtractedSource:
    return PARSER_REGISTRY.extract_source(path)


def _iter_source_files(root: Path) -> list[Path]:
    files: list[Path] = []
    for current, directories, names in os.walk(root):
        directories[:] = [
            name
            for name in directories
            if name not in IGNORED_DIRECTORY_NAMES and not name.startswith(".")
        ]
        base = Path(current)
        for name in names:
            path = base / name
            if path.is_symlink() or name.startswith("~$"):
                continue
            files.append(path)
    return sorted(files, key=lambda item: item.as_posix().casefold())


def _modified_at(stat: os.stat_result) -> str:
    return datetime.fromtimestamp(stat.st_mtime, tz=UTC).isoformat()


def _json_list(value: object) -> list[str]:
    if isinstance(value, list):
        return [str(item) for item in value]
    try:
        parsed = json.loads(str(value or "[]"))
    except (TypeError, ValueError):
        return []
    return [str(item) for item in parsed] if isinstance(parsed, list) else []


def _source_ref_json(source_ref: SourceRef) -> str:
    return json.dumps(source_ref.model_dump(mode="json"), ensure_ascii=False, separators=(",", ":"))


def _row_value(row: sqlite3.Row, key: str, default: object = "") -> object:
    columns = set(row.keys())
    return row[key] if key in columns else default


def _source_ref_from_row(row: sqlite3.Row) -> SourceRef:
    raw = str(_row_value(row, "source_ref_json") or "")
    if raw:
        try:
            return SourceRef.model_validate_json(raw)
        except ValueError:
            pass
    relative = str(row["relative_path"])
    kind = str(_row_value(row, "source_kind", "physical") or "physical")
    if kind == "archive_member":
        container = str(row["container_path"] or "")
        member = str(row["member_path"] or Path(relative).name)
        return SourceRef(
            kind="archive_member",
            workspace_path=container,
            virtual_path=relative,
            container_path=container,
            member_path=member,
            member_chain=[member],
            archive_depth=1,
            stable_id=str(row["document_id"]),
        )
    return physical_source_ref(relative, archive=kind == "archive").model_copy(
        update={"stable_id": str(row["document_id"])}
    )


def _quality_flags(value: object) -> list[str]:
    if isinstance(value, list):
        return sorted({str(item) for item in value if str(item)})
    try:
        parsed = json.loads(str(value or "[]"))
    except (TypeError, ValueError):
        return []
    return sorted({str(item) for item in parsed if str(item)}) if isinstance(parsed, list) else []


def _locator_from_json(value: object) -> EvidenceLocator | None:
    raw = str(value or "")
    if not raw:
        return None
    try:
        return EvidenceLocator.model_validate_json(raw)
    except ValueError:
        return None


def _freshness_from_value(value: object) -> Literal["current", "stale"]:
    return "stale" if str(value or "current").casefold() == "stale" else "current"


def _stable_snapshot(path: Path, *, attempts: int = 3) -> tuple[os.stat_result, str]:
    last_stat: os.stat_result | None = None
    last_hash = ""
    for _ in range(max(1, attempts)):
        first = path.stat()
        digest = sha256_file(path)
        second = path.stat()
        if first.st_size == second.st_size and first.st_mtime_ns == second.st_mtime_ns and digest:
            return second, digest
        last_stat, last_hash = second, digest
        time.sleep(0.05)
    if last_stat is None:
        raise FileNotFoundError(path)
    return last_stat, last_hash


def _passage_chunks(text: str, *, target: int = 1_600, overlap: int = 180) -> list[str]:
    compact = text.strip()
    if not compact:
        return []
    if len(compact) <= target:
        return [compact]
    chunks: list[str] = []
    start = 0
    while start < len(compact):
        desired_end = min(len(compact), start + target)
        end = desired_end
        if desired_end < len(compact):
            candidates = [
                compact.rfind(marker, start + target // 2, desired_end)
                for marker in ("\n\n", "\n", "。", ". ", "；", "; ")
            ]
            boundary = max(candidates, default=-1)
            if boundary > start:
                end = boundary + 1
        chunk = compact[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= len(compact):
            break
        start = max(start + 1, end - overlap)
    return chunks


def _apply_assisted_order(
    results: list[WorkspaceSearchResult],
    ordered_document_ids: list[str],
    exact_document_ids: list[str],
) -> list[WorkspaceSearchResult]:
    by_id = {result.document_id: result for result in results}
    allowed = set(by_id)
    exact = [document_id for document_id in exact_document_ids if document_id in allowed]
    requested: list[str] = []
    for document_id in ordered_document_ids:
        if document_id in allowed and document_id not in exact and document_id not in requested:
            requested.append(document_id)
    remaining = [
        result.document_id
        for result in results
        if result.document_id not in exact and result.document_id not in requested
    ]
    final_ids = [*exact, *requested, *remaining]
    return [
        by_id[document_id].model_copy(update={"rank": rank})
        for rank, document_id in enumerate(final_ids, start=1)
    ]


def assisted_rerank(
    workspace: GlobalWorkspace,
    query: str,
    results: list[WorkspaceSearchResult],
) -> tuple[list[WorkspaceSearchResult], str]:
    if not workspace.ai_policy.enabled or not results:
        raise RuntimeError("assisted_search_not_configured")
    try:
        credential = resolve_ai_api_key(workspace.workspace_id, workspace.ai_policy.provider)
    except CredentialStoreError as error:
        raise RuntimeError("credential_store_unavailable") from error
    if not credential.api_key:
        raise RuntimeError("ai_key_not_configured")
    from openai import OpenAI

    candidates = [
        {
            "document_id": result.document_id,
            "name": result.name,
            "relative_path": result.relative_path,
            "readability": result.readability,
            "reason": result.best_evidence.reason,
            "page_number": result.best_evidence.page_number,
            "excerpt": result.best_evidence.excerpt,
        }
        for result in results
    ]
    client = OpenAI(
        api_key=credential.api_key,
        base_url=workspace.ai_policy.base_url,
        max_retries=0,
        timeout=60.0,
    )
    response = client.chat.completions.create(
        model=workspace.ai_policy.model,
        messages=[
            {
                "role": "system",
                "content": (
                    "Reorder only the supplied document_id values for relevance. "
                    "Do not invent IDs or evidence. Return JSON with ordered_document_ids "
                    "and a concise Chinese answer grounded only in the candidates."
                ),
            },
            {
                "role": "user",
                "content": json.dumps(
                    {"query": query, "candidates": candidates}, ensure_ascii=False
                ),
            },
        ],
        response_format={"type": "json_object"},
        max_tokens=min(800, workspace.ai_policy.max_output_tokens_per_request),
        stream=False,
    )
    choices = getattr(response, "choices", None) or []
    content = getattr(getattr(choices[0], "message", None), "content", "") if choices else ""
    if not content:
        raise RuntimeError("ai_invalid_output")
    try:
        output = AssistedSearchOrder.model_validate_json(content)
    except ValueError as error:
        raise RuntimeError("ai_invalid_output") from error
    normalized_query = normalize_search_text(query)
    exact_ids = [
        result.document_id
        for result in results
        if normalized_query
        in {
            normalize_search_text(result.name),
            normalize_search_text(Path(result.name).stem),
        }
    ]
    return (
        _apply_assisted_order(results, output.ordered_document_ids, exact_ids),
        output.answer.strip(),
    )


def ensure_v2_workspaces() -> dict[str, GlobalWorkspace]:
    with global_config_lock():
        return _ensure_v2_workspaces_locked()


def _ensure_v2_workspaces_locked() -> dict[str, GlobalWorkspace]:
    config = load_global_config()
    changed = False
    for repository_id, repository in config.repositories.items():
        if repository_id in config.workspaces:
            continue
        legacy_path = Path(repository.index_repository_path)
        try:
            legacy_config = load_repository_config(legacy_path)
        except (OSError, ValueError):
            continue
        workspace = GlobalWorkspace(
            workspace_id=repository_id,
            name=repository.name,
            raw_path=legacy_config.repository.raw_repository_path,
            storage_path=str(workspace_storage_path(repository_id)),
            legacy_index_path=str(legacy_path),
            enabled=repository.enabled,
            ai_policy=legacy_config.ai_policy.model_copy(deep=True),
        )
        config.workspaces[repository_id] = workspace
        changed = True
    if not config.active_workspace_id and config.active_repository_id in config.workspaces:
        config.active_workspace_id = config.active_repository_id
        changed = True
    if changed:
        config.schema_version = WORKSPACE_SCHEMA_VERSION
        save_global_config(config)
    return config.workspaces


def create_workspace(raw_path: Path, name: str | None = None) -> GlobalWorkspace:
    raw = raw_path.expanduser().resolve()
    if not raw.is_dir():
        raise ValueError(f"资料文件夹不存在或不可访问: {raw}")
    with global_config_lock():
        _ensure_v2_workspaces_locked()
        config = load_global_config()
        for workspace in config.workspaces.values():
            existing_raw = Path(workspace.raw_path).expanduser().resolve()
            if existing_raw == raw:
                config.active_workspace_id = workspace.workspace_id
                save_global_config(config)
                return workspace
            if raw in existing_raw.parents or existing_raw in raw.parents:
                raise ValueError(
                    "Selected folder overlaps existing workspace "
                    f'"{workspace.name}" ({existing_raw}). Choose a non-overlapping folder '
                    "or use the existing workspace."
                )
        workspace_id = str(uuid.uuid4())
        workspace = GlobalWorkspace(
            workspace_id=workspace_id,
            name=(name or raw.name or "资料空间").strip(),
            raw_path=str(raw),
            storage_path=str(workspace_storage_path(workspace_id)),
            ai_policy=AIConfig(enabled=False),
        )
        config.schema_version = WORKSPACE_SCHEMA_VERSION
        config.workspaces[workspace_id] = workspace
        config.active_workspace_id = workspace_id
        save_global_config(config)
        return workspace


def get_workspace(workspace_id: str) -> GlobalWorkspace:
    ensure_v2_workspaces()
    workspace = load_global_config().workspaces.get(workspace_id)
    if workspace is None:
        raise KeyError(workspace_id)
    return workspace


class WorkspaceStore:
    def __init__(self, workspace: GlobalWorkspace) -> None:
        self.workspace = workspace
        self.raw = Path(workspace.raw_path).expanduser().resolve()
        self.storage = Path(workspace.storage_path).expanduser().resolve()
        self.database = self.storage / "workspace.sqlite3"
        self.previews = self.storage / "previews"

    def _connect(self) -> sqlite3.Connection:
        self.storage.mkdir(parents=True, exist_ok=True)
        connection = sqlite3.connect(self.database)
        connection.row_factory = sqlite3.Row
        connection.execute("PRAGMA foreign_keys = ON")
        connection.execute("PRAGMA journal_mode = WAL")
        self._ensure_schema(connection)
        return connection

    @staticmethod
    def _ensure_schema(connection: sqlite3.Connection) -> None:
        connection.executescript(
            """
            CREATE TABLE IF NOT EXISTS workspace_metadata (
                key TEXT PRIMARY KEY,
                value TEXT NOT NULL
            );
            CREATE TABLE IF NOT EXISTS documents (
                document_id TEXT PRIMARY KEY,
                relative_path TEXT NOT NULL UNIQUE,
                name TEXT NOT NULL,
                extension TEXT NOT NULL,
                content_hash TEXT NOT NULL,
                size_bytes INTEGER NOT NULL,
                mtime_ns INTEGER NOT NULL DEFAULT 0,
                modified_at TEXT NOT NULL,
                title TEXT NOT NULL,
                overview TEXT NOT NULL,
                page_count INTEGER NOT NULL,
                readability TEXT NOT NULL,
                readability_score REAL NOT NULL,
                indexing_state TEXT NOT NULL,
                error TEXT NOT NULL,
                quality_flags_json TEXT NOT NULL DEFAULT '[]',
                error_code TEXT NOT NULL DEFAULT '',
                parser_key TEXT NOT NULL DEFAULT '',
                parser_version TEXT NOT NULL DEFAULT '',
                source_kind TEXT NOT NULL DEFAULT 'physical',
                container_path TEXT NOT NULL DEFAULT '',
                member_path TEXT NOT NULL DEFAULT '',
                source_ref_json TEXT NOT NULL DEFAULT '',
                parent_document_id TEXT NOT NULL DEFAULT '',
                freshness_status TEXT NOT NULL DEFAULT 'current',
                indexed_at TEXT NOT NULL
            );
            CREATE INDEX IF NOT EXISTS documents_content_hash ON documents(content_hash);
            CREATE TABLE IF NOT EXISTS pages (
                page_id INTEGER PRIMARY KEY AUTOINCREMENT,
                document_id TEXT NOT NULL REFERENCES documents(document_id) ON DELETE CASCADE,
                page_number INTEGER,
                text TEXT NOT NULL,
                extraction_method TEXT NOT NULL,
                quality_score REAL NOT NULL,
                readability TEXT NOT NULL,
                locator_json TEXT NOT NULL DEFAULT ''
            );
            CREATE TABLE IF NOT EXISTS passages (
                passage_id INTEGER PRIMARY KEY AUTOINCREMENT,
                document_id TEXT NOT NULL REFERENCES documents(document_id) ON DELETE CASCADE,
                page_number INTEGER,
                ordinal INTEGER NOT NULL,
                heading TEXT NOT NULL,
                text TEXT NOT NULL,
                locator_json TEXT NOT NULL DEFAULT ''
            );
            CREATE VIRTUAL TABLE IF NOT EXISTS passages_fts USING fts5(
                document_id UNINDEXED,
                page_number UNINDEXED,
                name,
                path,
                heading,
                body,
                tokens,
                tokenize='unicode61 remove_diacritics 2'
            );
            CREATE TABLE IF NOT EXISTS change_events (
                change_id TEXT PRIMARY KEY,
                kind TEXT NOT NULL,
                document_id TEXT NOT NULL DEFAULT '',
                name TEXT NOT NULL DEFAULT '',
                relative_path TEXT NOT NULL DEFAULT '',
                previous_path TEXT NOT NULL DEFAULT '',
                occurred_at TEXT NOT NULL,
                message TEXT NOT NULL DEFAULT '',
                affected_task_ids_json TEXT NOT NULL DEFAULT '[]',
                acknowledged INTEGER NOT NULL DEFAULT 0
            );
            CREATE INDEX IF NOT EXISTS change_events_occurred_at
                ON change_events(occurred_at DESC);
            """
        )
        document_columns = {
            str(row[1]) for row in connection.execute("PRAGMA table_info(documents)").fetchall()
        }
        if "mtime_ns" not in document_columns:
            connection.execute(
                "ALTER TABLE documents ADD COLUMN mtime_ns INTEGER NOT NULL DEFAULT 0"
            )
        additive_document_columns = {
            "quality_flags_json": "TEXT NOT NULL DEFAULT '[]'",
            "error_code": "TEXT NOT NULL DEFAULT ''",
            "parser_key": "TEXT NOT NULL DEFAULT ''",
            "parser_version": "TEXT NOT NULL DEFAULT ''",
            "source_kind": "TEXT NOT NULL DEFAULT 'physical'",
            "container_path": "TEXT NOT NULL DEFAULT ''",
            "member_path": "TEXT NOT NULL DEFAULT ''",
            "source_ref_json": "TEXT NOT NULL DEFAULT ''",
            "parent_document_id": "TEXT NOT NULL DEFAULT ''",
            "freshness_status": "TEXT NOT NULL DEFAULT 'current'",
        }
        for name, declaration in additive_document_columns.items():
            if name not in document_columns:
                connection.execute(f"ALTER TABLE documents ADD COLUMN {name} {declaration}")
        for table in ("pages", "passages"):
            columns = {
                str(row[1]) for row in connection.execute(f"PRAGMA table_info({table})").fetchall()
            }
            if "locator_json" not in columns:
                connection.execute(
                    f"ALTER TABLE {table} ADD COLUMN locator_json TEXT NOT NULL DEFAULT ''"
                )
        connection.execute(
            "CREATE INDEX IF NOT EXISTS documents_source_kind ON documents(source_kind)"
        )
        connection.execute(
            "CREATE INDEX IF NOT EXISTS documents_container_path ON documents(container_path)"
        )

    def _source_path(self, relative_path: str) -> Path:
        path = (self.raw / relative_path).resolve()
        if path != self.raw and self.raw not in path.parents:
            raise ValueError("Source path escapes the workspace")
        return path

    def _archive_policy(self) -> ArchivePolicy:
        policy = self.workspace.archive_policy
        return ArchivePolicy(
            max_members=policy.max_entries,
            max_member_bytes=policy.max_member_bytes,
            max_total_bytes=policy.max_total_bytes,
            max_compression_ratio=policy.max_compression_ratio,
            max_nested_archives=policy.nested_zip_depth,
            cache_ttl_seconds=policy.materialized_cache_ttl_hours * 60 * 60,
            cache_max_bytes=policy.materialized_cache_max_bytes,
        )

    def _member_cache_root(self) -> Path:
        return self.storage / "member-cache"

    def _affected_task_ids(
        self,
        document_id: str,
        paths: Sequence[str],
    ) -> list[str]:
        """Find task packages that reference a changed physical or virtual path."""
        task_directory = workspace_tasks_path(self.workspace.workspace_id)
        if not task_directory.is_dir():
            return []
        normalized_paths = {
            value.replace("\\", "/").casefold().strip("/")
            for value in paths
            if value
        }
        affected: set[str] = set()
        for task_path_value in task_directory.glob("*.json"):
            payload = load_json(task_path_value, {})
            if not isinstance(payload, dict):
                continue
            raw_items = payload.get("items", [])
            if not isinstance(raw_items, list):
                continue
            for raw_item in raw_items:
                if not isinstance(raw_item, dict):
                    continue
                if document_id and str(raw_item.get("document_id", "")) == document_id:
                    affected.add(str(payload.get("task_id") or task_path_value.stem))
                    break
                item_path = str(raw_item.get("relative_path", ""))
                source_ref = raw_item.get("source_ref")
                if isinstance(source_ref, dict):
                    item_path = str(source_ref.get("virtual_path") or item_path)
                if item_path.replace("\\", "/").casefold().strip("/") in normalized_paths:
                    affected.add(str(payload.get("task_id") or task_path_value.stem))
                    break
        return sorted(affected)

    def _record_change(
        self,
        connection: sqlite3.Connection,
        *,
        kind: Literal["added", "modified", "moved", "deleted", "parser_warning"],
        document_id: str = "",
        name: str = "",
        relative_path: str = "",
        previous_path: str = "",
        message: str = "",
    ) -> None:
        paths = [relative_path, previous_path]
        affected = self._affected_task_ids(document_id, paths)
        connection.execute(
            "INSERT INTO change_events("
            "change_id, kind, document_id, name, relative_path, previous_path, "
            "occurred_at, message, affected_task_ids_json, acknowledged) "
            "VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, 0)",
            (
                uuid.uuid4().hex,
                kind,
                document_id,
                name,
                relative_path,
                previous_path,
                utc_now(),
                message[:500],
                json.dumps(affected, ensure_ascii=False),
            ),
        )

    def _source_uri(self, row: sqlite3.Row) -> str:
        source_ref = _source_ref_from_row(row)
        relative = (
            source_ref.container_path
            if source_ref.source_kind == "archive_member" and source_ref.container_path
            else str(row["relative_path"])
        )
        try:
            return self._source_path(relative).as_uri()
        except ValueError:
            return ""

    @staticmethod
    def _archive_document_id(
        workspace_id: str,
        source_ref: SourceRef,
        container_id: str = "",
    ) -> str:
        identity = "\x00".join(
            [
                workspace_id,
                container_id or source_ref.container_path,
                "/".join(source_ref.member_chain),
                ",".join(str(item) for item in source_ref.member_indexes),
            ]
        )
        return str(uuid.uuid5(uuid.NAMESPACE_URL, identity))

    def _archive_extracted(
        self,
        candidate: ArchiveCandidate,
    ) -> ExtractedSource:
        if candidate.materialized_path is None or not candidate.materialized_path.is_file():
            return ExtractedSource(
                title=Path(candidate.display_name).stem,
                pages=[],
                page_count=0,
                status="metadata_only",
                error=candidate.error_code,
                quality_flags=sorted(set(candidate.quality_flags)),
                error_code=candidate.error_code,
                parser_key="v2.archive_member",
                parser_version="1",
            )
        extracted = extract_source(candidate.materialized_path)
        flags = sorted(set([*extracted.quality_flags, *candidate.quality_flags]))
        error_code = extracted.error_code or candidate.error_code
        return replace(
            extracted,
            quality_flags=flags,
            error_code=error_code,
        )

    def _sync_archive_members(
        self,
        connection: sqlite3.Connection,
        *,
        path: Path,
        relative: str,
        container_id: str,
        container_hash: str,
        outer_stat: os.stat_result,
        existing: dict[str, sqlite3.Row],
        seen: set[str],
        report_progress: Callable[..., None],
        force_document_ids: set[str],
    ) -> tuple[int, int]:
        if not self.workspace.archive_policy.enabled:
            return 0, 0
        cache_root = self._member_cache_root()
        scan = scan_archive(
            path,
            root_relative=relative,
            cache_root=cache_root,
            policy=self._archive_policy(),
        )
        if scan.quality_flags or scan.error_code:
            connection.execute(
                "UPDATE documents SET quality_flags_json = ?, error_code = ?, "
                "error = ? WHERE document_id = ?",
                (
                    json.dumps(
                        sorted(set(["archive_container", *scan.quality_flags])), ensure_ascii=False
                    ),
                    scan.error_code,
                    scan.error[:500],
                    container_id,
                ),
            )
            self._record_change(
                connection,
                kind="parser_warning",
                document_id=container_id,
                name=path.name,
                relative_path=relative,
                message=scan.error_code or ", ".join(scan.quality_flags),
            )
        if scan.error_code and not scan.members:
            # Keep the last usable member index. It is stale, but deleting it would
            # make a transiently corrupt/syncing archive look like data loss.
            for old_path, old_row in existing.items():
                if str(old_row["container_path"] or "") == relative:
                    seen.add(old_path)
                    connection.execute(
                        "UPDATE documents SET freshness_status = 'stale' WHERE document_id = ?",
                        (str(old_row["document_id"]),),
                    )
            return 0, len(scan.members)

        indexed = 0
        failed = 0
        for member_number, candidate in enumerate(scan.members, start=1):
            virtual_path = candidate.virtual_path
            seen.add(virtual_path)
            report_progress(
                current_archive=relative,
                current_member=virtual_path,
                member_processed=member_number,
                member_total=len(scan.members),
            )
            current = existing.get(virtual_path)
            content_hash = (
                candidate.content_hash
                or hashlib.sha256(
                    f"{container_hash}:{virtual_path}:{candidate.size_bytes}".encode()
                ).hexdigest()
            )
            parser_key, parser_version = (
                parser_signature(candidate.materialized_path)
                if candidate.materialized_path is not None
                else ("v2.archive_member", "1")
            )
            if (
                current is not None
                and str(current["document_id"]) not in force_document_ids
                and str(current["content_hash"]) == content_hash
                and str(current["parser_key"]) == parser_key
                and str(current["parser_version"]) == parser_version
                and str(current["indexing_state"]) != "failed"
            ):
                continue
            document_id = str(current["document_id"]) if current is not None else ""
            moved_row: sqlite3.Row | None = None
            if not document_id:
                candidate_chain = tuple(candidate.source_ref.member_chain)
                candidate_indexes = tuple(candidate.source_ref.member_indexes)
                matches = []
                for row in existing.values():
                    if str(row["parent_document_id"] or "") != container_id:
                        continue
                    if str(row["content_hash"]) != content_hash:
                        continue
                    old_ref = _source_ref_from_row(row)
                    if (
                        tuple(old_ref.member_chain) == candidate_chain
                        and tuple(old_ref.member_indexes) == candidate_indexes
                    ):
                        matches.append(row)
                moved_row = matches[0] if len(matches) == 1 else None
                document_id = (
                    self._archive_document_id(
                        self.workspace.workspace_id,
                        candidate.source_ref,
                        container_id,
                    )
                    if moved_row is None
                    else str(moved_row["document_id"])
                )
            source_ref = candidate.source_ref.model_copy(update={"stable_id": document_id})
            extracted = self._archive_extracted(candidate)
            try:
                connection.execute("SAVEPOINT replace_archive_member")
                self._replace_document(
                    connection,
                    document_id=document_id,
                    relative_path=virtual_path,
                    path=candidate.materialized_path or path,
                    name=candidate.display_name,
                    extension=candidate.extension,
                    content_hash=content_hash,
                    size_bytes=candidate.size_bytes,
                    mtime_ns=outer_stat.st_mtime_ns,
                    modified_at=candidate.modified_at,
                    extracted=extracted,
                    source_ref=source_ref,
                    parent_document_id=container_id,
                )
                connection.execute("RELEASE SAVEPOINT replace_archive_member")
                if moved_row is not None:
                    self._record_change(
                        connection,
                        kind="moved",
                        document_id=document_id,
                        name=candidate.display_name,
                        relative_path=virtual_path,
                        previous_path=str(moved_row["relative_path"]),
                    )
                elif current is None:
                    self._record_change(
                        connection,
                        kind="added",
                        document_id=document_id,
                        name=candidate.display_name,
                        relative_path=virtual_path,
                    )
                else:
                    self._record_change(
                        connection,
                        kind="modified",
                        document_id=document_id,
                        name=candidate.display_name,
                        relative_path=virtual_path,
                    )
                indexed += 1
            except Exception as error:
                try:
                    connection.execute("ROLLBACK TO SAVEPOINT replace_archive_member")
                    connection.execute("RELEASE SAVEPOINT replace_archive_member")
                except sqlite3.OperationalError:
                    pass
                self._replace_failed_document(
                    connection,
                    document_id=document_id,
                    relative_path=virtual_path,
                    path=candidate.materialized_path or path,
                    name=candidate.display_name,
                    extension=candidate.extension,
                    content_hash=content_hash,
                    size_bytes=candidate.size_bytes,
                    mtime_ns=outer_stat.st_mtime_ns,
                    modified_at=candidate.modified_at,
                    error=error,
                    source_ref=source_ref,
                    parent_document_id=container_id,
                )
                self._record_change(
                    connection,
                    kind="parser_warning",
                    document_id=document_id,
                    name=candidate.display_name,
                    relative_path=virtual_path,
                    message=f"{type(error).__name__}: {error}",
                )
                failed += 1
        return indexed, failed

    def sync(
        self,
        progress_callback: Callable[[dict[str, Any]], None] | None = None,
        *,
        force_document_ids: set[str] | None = None,
    ) -> dict[str, Any]:
        if not self.raw.is_dir():
            raise FileNotFoundError(f"资料文件夹不可访问: {self.raw}")
        progress_state: dict[str, Any] = {
            "phase": "discovering",
            "discovered": 0,
            "processed": 0,
            "current_file": "",
            "indexed": 0,
            "unchanged": 0,
            "failed": 0,
            "removed": 0,
        }

        def report_progress(**changes: Any) -> None:
            if changes.pop("clear_page_progress", False):
                for key in (
                    "current_page",
                    "page_count",
                    "pages_completed",
                    "ocr_pages_completed",
                    "extraction_stage",
                ):
                    progress_state.pop(key, None)
            progress_state.update(changes)
            if progress_callback is not None:
                progress_callback(dict(progress_state))

        report_progress()
        files = _iter_source_files(self.raw)
        report_progress(phase="processing", discovered=len(files))
        discovered_paths = {path.relative_to(self.raw).as_posix() for path in files}
        seen: set[str] = set()
        indexed = 0
        unchanged = 0
        failed = 0
        processed = 0
        forced = force_document_ids or set()
        with closing(self._connect()) as connection:
            existing = {
                str(row["relative_path"]): row
                for row in connection.execute("SELECT * FROM documents").fetchall()
            }
            movable_by_identity: dict[tuple[str, str], list[sqlite3.Row]] = {}
            for relative_path, row in existing.items():
                source_kind = str(_row_value(row, "source_kind", "physical") or "physical")
                if (
                    relative_path not in discovered_paths
                    and source_kind in {"physical", "archive"}
                    and not str(_row_value(row, "parent_document_id", "") or "")
                ):
                    movable_by_identity.setdefault(
                        (source_kind, str(row["content_hash"])), []
                    ).append(row)
            for path in files:
                relative = path.relative_to(self.raw).as_posix()
                report_progress(clear_page_progress=True, current_file=relative)
                seen.add(relative)
                stat, content_hash = _stable_snapshot(
                    path,
                    attempts=max(1, self.workspace.sync_policy.stable_retry_count),
                )
                current = existing.get(relative)
                expected_parser_key, expected_parser_version = parser_signature(path)
                if (
                    current is not None
                    and str(current["document_id"]) not in forced
                    and str(current["indexing_state"]) != "failed"
                    and int(current["size_bytes"]) == stat.st_size
                    and int(current["mtime_ns"]) == stat.st_mtime_ns
                    and str(current["content_hash"]) == content_hash
                    and str(current["parser_key"]) == expected_parser_key
                    and str(current["parser_version"]) == expected_parser_version
                ):
                    if path.suffix.casefold() == ".zip" and self.workspace.archive_policy.enabled:
                        for member_path, member_row in existing.items():
                            if (
                                str(_row_value(member_row, "source_kind", ""))
                                == "archive_member"
                                and str(_row_value(member_row, "container_path", ""))
                                == relative
                            ):
                                seen.add(member_path)
                    unchanged += 1
                    processed += 1
                    report_progress(
                        processed=processed,
                        indexed=indexed,
                        unchanged=unchanged,
                        failed=failed,
                    )
                    continue
                document_id = str(current["document_id"]) if current else ""
                moved_row: sqlite3.Row | None = None
                if not document_id:
                    source_kind = "archive" if path.suffix.casefold() == ".zip" else "physical"
                    moved_candidates = movable_by_identity.get((source_kind, content_hash), [])
                    moved_row = moved_candidates.pop(0) if moved_candidates else None
                    document_id = (
                        str(moved_row["document_id"])
                        if moved_row is not None
                        else str(uuid.uuid4())
                    )
                try:
                    progress_token = _EXTRACTION_PROGRESS_CALLBACK.set(
                        lambda update: report_progress(**update)
                    )
                    try:
                        extracted = extract_source(path)
                        source_ref = physical_source_ref(
                            relative,
                            archive=path.suffix.casefold() == ".zip",
                        )
                        if path.suffix.casefold() == ".zip":
                            extracted = ExtractedSource(
                                title=path.stem,
                                pages=[],
                                page_count=0,
                                status="metadata_only",
                                quality_flags=["archive_container"],
                                parser_key="v2.archive",
                                parser_version="1",
                            )
                    finally:
                        _EXTRACTION_PROGRESS_CALLBACK.reset(progress_token)
                    connection.execute("SAVEPOINT replace_document")
                    self._replace_document(
                        connection,
                        document_id=document_id,
                        relative_path=relative,
                        path=path,
                        name=path.name,
                        extension=path.suffix.casefold(),
                        content_hash=content_hash,
                        size_bytes=stat.st_size,
                        mtime_ns=stat.st_mtime_ns,
                        modified_at=_modified_at(stat),
                        extracted=extracted,
                        source_ref=source_ref,
                    )
                    connection.execute("RELEASE SAVEPOINT replace_document")
                    indexed += 1
                    if moved_row is not None:
                        self._record_change(
                            connection,
                            kind="moved",
                            document_id=document_id,
                            name=path.name,
                            relative_path=relative,
                            previous_path=str(moved_row["relative_path"]),
                        )
                    elif current is None:
                        self._record_change(
                            connection,
                            kind="added",
                            document_id=document_id,
                            name=path.name,
                            relative_path=relative,
                        )
                    else:
                        self._record_change(
                            connection,
                            kind="modified",
                            document_id=document_id,
                            name=path.name,
                            relative_path=relative,
                        )
                    if path.suffix.casefold() == ".zip":
                        member_indexed, member_failed = self._sync_archive_members(
                            connection,
                            path=path,
                            relative=relative,
                            container_id=document_id,
                            container_hash=content_hash,
                            outer_stat=stat,
                            existing=existing,
                            seen=seen,
                            report_progress=report_progress,
                            force_document_ids=forced,
                        )
                        indexed += member_indexed
                        failed += member_failed
                except Exception as error:
                    try:
                        connection.execute("ROLLBACK TO SAVEPOINT replace_document")
                        connection.execute("RELEASE SAVEPOINT replace_document")
                    except sqlite3.OperationalError:
                        pass
                    self._replace_failed_document(
                        connection,
                        document_id=document_id,
                        relative_path=relative,
                        path=path,
                        name=path.name,
                        extension=path.suffix.casefold(),
                        content_hash=content_hash,
                        size_bytes=stat.st_size,
                        mtime_ns=stat.st_mtime_ns,
                        modified_at=_modified_at(stat),
                        error=error,
                        source_ref=physical_source_ref(
                            relative,
                            archive=path.suffix.casefold() == ".zip",
                        ),
                    )
                    self._record_change(
                        connection,
                        kind="parser_warning",
                        document_id=document_id,
                        name=path.name,
                        relative_path=relative,
                        message=f"{type(error).__name__}: {error}",
                    )
                    failed += 1
                processed += 1
                report_progress(
                    processed=processed,
                    indexed=indexed,
                    unchanged=unchanged,
                    failed=failed,
                )
            removed = sorted(set(existing) - seen)
            report_progress(phase="finalizing", current_file="", removed=len(removed))
            for relative in removed:
                document_id = str(existing[relative]["document_id"])
                current_path = connection.execute(
                    "SELECT relative_path FROM documents WHERE document_id = ?", (document_id,)
                ).fetchone()
                if current_path is not None and str(current_path["relative_path"]) != relative:
                    continue
                self._record_change(
                    connection,
                    kind="deleted",
                    document_id=document_id,
                    name=str(existing[relative]["name"]),
                    relative_path=relative,
                )
                self._delete_document(connection, document_id)
            now = utc_now()
            connection.execute(
                "INSERT INTO workspace_metadata(key, value) VALUES('last_sync_at', ?) "
                "ON CONFLICT(key) DO UPDATE SET value = excluded.value",
                (now,),
            )
            connection.execute(
                "INSERT INTO workspace_metadata(key, value) VALUES('schema_version', ?) "
                "ON CONFLICT(key) DO UPDATE SET value = excluded.value",
                (WORKSPACE_SCHEMA_VERSION,),
            )
            connection.commit()
        result = {
            "workspace_id": self.workspace.workspace_id,
            "discovered": len(files),
            "indexed": indexed,
            "unchanged": unchanged,
            "removed": len(removed),
            "failed": failed,
            "health": self.health().model_dump(mode="json"),
        }
        report_progress(phase="completed", current_file="")
        return result

    def _delete_document(self, connection: sqlite3.Connection, document_id: str) -> None:
        rows = connection.execute(
            "SELECT passage_id FROM passages WHERE document_id = ?", (document_id,)
        ).fetchall()
        for row in rows:
            connection.execute("DELETE FROM passages_fts WHERE rowid = ?", (int(row[0]),))
        connection.execute("DELETE FROM documents WHERE document_id = ?", (document_id,))

    def _replace_document(
        self,
        connection: sqlite3.Connection,
        *,
        document_id: str,
        relative_path: str,
        path: Path,
        name: str,
        extension: str,
        content_hash: str,
        size_bytes: int,
        mtime_ns: int,
        modified_at: str,
        extracted: ExtractedSource,
        source_ref: SourceRef,
        parent_document_id: str = "",
    ) -> None:
        if connection.execute(
            "SELECT 1 FROM documents WHERE document_id = ?", (document_id,)
        ).fetchone():
            self._delete_document(connection, document_id)
        readable_pages = [
            page for page in extracted.pages if page.quality_score >= READABLE_THRESHOLD
        ]
        scores = [page.quality_score for page in extracted.pages]
        document_score = sum(scores) / len(scores) if scores else 0.0
        readability = readability_label(document_score)
        if extracted.status == "metadata_only":
            readability = "low"
        overview = ""
        if extracted.pages and len(readable_pages) / len(extracted.pages) >= 0.70:
            overview = _excerpt(readable_pages[0].text, "", 260)
        stable_source_ref = source_ref.model_copy(update={"stable_id": document_id})
        parser_key = extracted.parser_key or parser_signature(path)[0]
        parser_version = extracted.parser_version or parser_signature(path)[1]
        title = extracted.title
        if source_ref.source_kind == "archive_member" and title == path.stem:
            title = Path(name).stem
        connection.execute(
            """
            INSERT INTO documents(
                document_id, relative_path, name, extension, content_hash, size_bytes, mtime_ns,
                modified_at, title, overview, page_count, readability, readability_score,
                indexing_state, error, quality_flags_json, error_code, parser_key, parser_version,
                source_kind, container_path, member_path, source_ref_json, parent_document_id,
                freshness_status, indexed_at
            ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?,
                      ?, ?)
            """,
            (
                document_id,
                relative_path,
                name,
                extension,
                content_hash,
                size_bytes,
                mtime_ns,
                modified_at,
                title,
                overview,
                extracted.page_count,
                readability,
                round(document_score, 4),
                extracted.status,
                extracted.error,
                json.dumps(sorted(set(extracted.quality_flags)), ensure_ascii=False),
                extracted.error_code,
                parser_key,
                parser_version,
                source_ref.source_kind,
                source_ref.container_path,
                source_ref.member_path,
                _source_ref_json(stable_source_ref),
                parent_document_id,
                "current",
                utc_now(),
            ),
        )
        for page in extracted.pages:
            locator = page.locator
            if locator is None and page.page_number is not None:
                locator = EvidenceLocator(kind="page", page_number=page.page_number)
            locator_json = (
                json.dumps(locator.model_dump(mode="json"), ensure_ascii=False)
                if locator is not None
                else ""
            )
            connection.execute(
                "INSERT INTO pages("
                "document_id, page_number, text, extraction_method, quality_score, readability, "
                "locator_json"
                ") "
                "VALUES (?, ?, ?, ?, ?, ?, ?)",
                (
                    document_id,
                    page.page_number,
                    page.text,
                    page.extraction_method,
                    page.quality_score,
                    page.readability,
                    locator_json,
                ),
            )
            if page.quality_score < PARTIAL_THRESHOLD or not page.text.strip():
                continue
            heading = next((line.strip() for line in page.text.splitlines() if line.strip()), "")[
                :200
            ]
            for ordinal, chunk in enumerate(_passage_chunks(page.text)):
                passage = connection.execute(
                    "INSERT INTO passages("
                    "document_id, page_number, ordinal, heading, text, locator_json"
                    ") VALUES (?, ?, ?, ?, ?, ?)",
                    (document_id, page.page_number, ordinal, heading, chunk, locator_json),
                )
                if passage.lastrowid is None:
                    raise RuntimeError("Unable to allocate passage ID")
                passage_id = passage.lastrowid
                connection.execute(
                    "INSERT INTO passages_fts("
                    "rowid, document_id, page_number, name, path, heading, body, tokens"
                    ") VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
                    (
                        passage_id,
                        document_id,
                        page.page_number,
                        name,
                        relative_path,
                        heading,
                        chunk,
                        _fts_text(" ".join([path.name, relative_path, heading, chunk])),
                    ),
                )

    def _replace_failed_document(
        self,
        connection: sqlite3.Connection,
        *,
        document_id: str,
        relative_path: str,
        path: Path,
        name: str,
        extension: str,
        content_hash: str,
        size_bytes: int,
        mtime_ns: int,
        modified_at: str,
        error: Exception,
        source_ref: SourceRef,
        parent_document_id: str = "",
    ) -> None:
        if connection.execute(
            "SELECT 1 FROM documents WHERE document_id = ?", (document_id,)
        ).fetchone():
            self._delete_document(connection, document_id)
        parser_key, parser_version = parser_signature(path)
        stable_source_ref = source_ref.model_copy(update={"stable_id": document_id})
        connection.execute(
            """
            INSERT INTO documents(
                document_id, relative_path, name, extension, content_hash, size_bytes, mtime_ns,
                modified_at, title, overview, page_count, readability, readability_score,
                indexing_state, error, quality_flags_json, error_code, parser_key, parser_version,
                source_kind, container_path, member_path, source_ref_json, parent_document_id,
                freshness_status, indexed_at
            ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, '', 0, 'low', 0, 'failed', ?, '[]',
                      'parser_failed', ?, ?, ?, ?, ?, ?, ?, 'current', ?)
            """,
            (
                document_id,
                relative_path,
                name,
                extension,
                content_hash,
                size_bytes,
                mtime_ns,
                modified_at,
                Path(name).stem,
                f"{type(error).__name__}: {error}"[:500],
                parser_key,
                parser_version,
                source_ref.source_kind,
                source_ref.container_path,
                source_ref.member_path,
                _source_ref_json(stable_source_ref),
                parent_document_id,
                utc_now(),
            ),
        )

    def health(self) -> WorkspaceHealth:
        if not self.database.exists():
            return WorkspaceHealth()
        with closing(self._connect()) as connection:
            document_count = int(connection.execute("SELECT COUNT(*) FROM documents").fetchone()[0])
            counts = {
                str(row["readability"]): int(row["count"])
                for row in connection.execute(
                    "SELECT readability, COUNT(*) AS count FROM documents "
                    "WHERE indexing_state = 'indexed' GROUP BY readability"
                ).fetchall()
            }
            metadata_only = int(
                connection.execute(
                    "SELECT COUNT(*) FROM documents WHERE indexing_state = 'metadata_only'"
                ).fetchone()[0]
            )
            failed = int(
                connection.execute(
                    "SELECT COUNT(*) FROM documents WHERE indexing_state = 'failed'"
                ).fetchone()[0]
            )
            last_sync = connection.execute(
                "SELECT value FROM workspace_metadata WHERE key = 'last_sync_at'"
            ).fetchone()
            return WorkspaceHealth(
                document_count=document_count,
                readable_count=counts.get("readable", 0),
                partial_count=counts.get("partial", 0),
                low_quality_count=counts.get("low", 0),
                metadata_only_count=metadata_only,
                failed_count=failed,
                last_sync_at=str(last_sync[0]) if last_sync else "",
            )

    def payload(self) -> WorkspacePayload:
        return WorkspacePayload(
            workspace_id=self.workspace.workspace_id,
            name=self.workspace.name,
            raw_path=str(self.raw),
            available=self.raw.is_dir(),
            enabled=self.workspace.enabled,
            vision_enabled=self.workspace.vision_enabled,
            legacy_index_present=bool(
                self.workspace.legacy_index_path and Path(self.workspace.legacy_index_path).exists()
            ),
            health=self.health(),
        )

    def list_documents(self) -> list[WorkspaceDocument]:
        if not self.database.exists():
            return []
        with closing(self._connect()) as connection:
            rows = connection.execute(
                "SELECT * FROM documents ORDER BY relative_path COLLATE NOCASE"
            ).fetchall()
        return [self._document_payload(row) for row in rows]

    def get_document(self, document_id: str) -> WorkspaceDocument:
        with closing(self._connect()) as connection:
            row = connection.execute(
                "SELECT * FROM documents WHERE document_id = ?", (document_id,)
            ).fetchone()
        if row is None:
            raise FileNotFoundError("Document not found")
        return self._document_payload(row)

    def list_members(self, container_document_id: str) -> list[WorkspaceDocument]:
        with closing(self._connect()) as connection:
            container = connection.execute(
                "SELECT relative_path FROM documents WHERE document_id = ?",
                (container_document_id,),
            ).fetchone()
            if container is None:
                raise FileNotFoundError("Document not found")
            rows = connection.execute(
                "SELECT * FROM documents WHERE source_kind = 'archive_member' "
                "AND (parent_document_id = ? OR container_path = ?) "
                "ORDER BY relative_path COLLATE NOCASE",
                (container_document_id, str(container["relative_path"])),
            ).fetchall()
        return [self._document_payload(row) for row in rows]

    def list_changes(
        self,
        *,
        limit: int = 100,
        since: str = "",
        include_acknowledged: bool = False,
    ) -> list[dict[str, Any]]:
        """Return recent source changes without loading the whole index into Python."""
        bounded_limit = max(1, min(int(limit), 1_000))
        clauses: list[str] = []
        parameters: list[object] = []
        if since.strip():
            clauses.append("occurred_at > ?")
            parameters.append(since.strip())
        if not include_acknowledged:
            clauses.append("acknowledged = 0")
        where = f"WHERE {' AND '.join(clauses)}" if clauses else ""
        with closing(self._connect()) as connection:
            rows = connection.execute(
                "SELECT change_id, kind, document_id, name, relative_path, "
                "previous_path, occurred_at, message, affected_task_ids_json, acknowledged "
                f"FROM change_events {where} ORDER BY occurred_at DESC LIMIT ?",
                [*parameters, bounded_limit],
            ).fetchall()
        changes: list[dict[str, Any]] = []
        for row in rows:
            try:
                affected = json.loads(str(row["affected_task_ids_json"] or "[]"))
            except (TypeError, ValueError):
                affected = []
            changes.append(
                {
                    "change_id": str(row["change_id"]),
                    "kind": str(row["kind"]),
                    "document_id": str(row["document_id"]),
                    "name": str(row["name"]),
                    "relative_path": str(row["relative_path"]),
                    "previous_path": str(row["previous_path"]),
                    "occurred_at": str(row["occurred_at"]),
                    "message": str(row["message"]),
                    "affected_task_ids": (
                        [str(item) for item in affected]
                        if isinstance(affected, list)
                        else []
                    ),
                    "acknowledged": bool(row["acknowledged"]),
                }
            )
        return changes

    def content_path(self, document_id: str) -> Path:
        with closing(self._connect()) as connection:
            row = connection.execute(
                "SELECT * FROM documents WHERE document_id = ?", (document_id,)
            ).fetchone()
        if row is None:
            raise FileNotFoundError("Document not found")
        source_ref = _source_ref_from_row(row)
        if source_ref.source_kind == "archive_member":
            return materialize_source_ref(
                self.raw,
                source_ref,
                cache_root=self._member_cache_root(),
                expected_hash=str(row["content_hash"]),
                policy=self._archive_policy(),
            )
        source = self._source_path(str(row["relative_path"]))
        if not source.is_file():
            raise FileNotFoundError("Source file is unavailable")
        return source

    def content_bytes(self, document_id: str, *, max_bytes: int = 100 * 1024 * 1024) -> bytes:
        path = self.content_path(document_id)
        if path.stat().st_size > max_bytes:
            raise ValueError("Document exceeds the content preview limit")
        return path.read_bytes()

    def materialize_document(self, document_id: str) -> Path:
        return self.content_path(document_id)

    def open_target(self, document_id: str) -> dict[str, Any]:
        path = self.content_path(document_id)
        document = self.get_document(document_id)
        temporary = bool(
            document.source_ref and document.source_ref.source_kind == "archive_member"
        )
        expires_at = cache_expiry(path, self._archive_policy()) if temporary else ""
        return {
            "uri": path.resolve().as_uri(),
            "temporary": temporary,
            "expires_at": expires_at,
            "display_name": document.name,
            "source_ref": document.source_ref.model_dump(mode="json")
            if document.source_ref
            else None,
        }

    def _document_payload(self, row: sqlite3.Row) -> WorkspaceDocument:
        source_ref = _source_ref_from_row(row)
        source_uri = self._source_uri(row)
        return WorkspaceDocument(
            document_id=str(row["document_id"]),
            name=str(row["name"]),
            relative_path=str(row["relative_path"]),
            extension=str(row["extension"]),
            content_hash=str(row["content_hash"]),
            size_bytes=int(row["size_bytes"]),
            modified_at=str(row["modified_at"]),
            title=str(row["title"]),
            overview=str(row["overview"]),
            page_count=int(row["page_count"]),
            readability=_readability_value(row["readability"]),
            readability_score=float(row["readability_score"]),
            indexing_state=_indexing_state_value(row["indexing_state"]),
            error=str(row["error"]),
            source_uri=source_uri,
            source_ref=source_ref,
            locator=_locator_from_json(_row_value(row, "locator_json")),
            quality_flags=_quality_flags(_row_value(row, "quality_flags_json", "[]")),
            error_code=str(_row_value(row, "error_code")),
            parser_key=str(_row_value(row, "parser_key")),
            parser_version=str(_row_value(row, "parser_version")),
            freshness_status=_freshness_from_value(
                _row_value(row, "freshness_status", "current")
            ),
        )

    def search(
        self,
        query: str,
        *,
        limit: int = 30,
        mode: Literal["local", "assisted"] = "local",
        path_prefix: str = "",
        extensions: Sequence[str] | None = None,
        readability: Sequence[str] | None = None,
        indexing_states: Sequence[str] | None = None,
        source_kinds: Sequence[str] | None = None,
        modified_from: str = "",
        modified_to: str = "",
        task_id: str = "",
    ) -> WorkspaceSearchReport:
        started = time.perf_counter()
        value = query.strip()
        if not value:
            raise ValueError("Search query cannot be empty")
        normalized_query = normalize_search_text(value)
        if not normalized_query:
            raise ValueError("Search query must contain searchable text")
        with closing(self._connect()) as connection:
            allowed_extensions = {item.casefold() for item in (extensions or [])}
            allowed_readability = {item.casefold() for item in (readability or [])}
            allowed_states = {item.casefold() for item in (indexing_states or [])}
            allowed_kinds = {item.casefold() for item in (source_kinds or [])}
            task_document_ids: set[str] | None = None
            if task_id:
                try:
                    normalized_task_id = str(uuid.UUID(task_id))
                except ValueError as error:
                    raise ValueError("task_id must be a valid UUID") from error
                task_document_ids = set()
                task_path_value = (
                    workspace_tasks_path(self.workspace.workspace_id)
                    / f"{normalized_task_id}.json"
                )
                task_payload = load_json(task_path_value, {})
                raw_items = task_payload.get("items", []) if isinstance(task_payload, dict) else []
                for item in raw_items if isinstance(raw_items, list) else []:
                    if isinstance(item, dict) and item.get("document_id"):
                        task_document_ids.add(str(item["document_id"]))

            filters: list[str] = []
            filter_parameters: list[object] = []

            def add_values_filter(column: str, values: set[str]) -> None:
                if not values:
                    return
                placeholders = ", ".join("?" for _ in values)
                filters.append(f"{column} IN ({placeholders})")
                filter_parameters.extend(sorted(values))

            if path_prefix:
                escaped_prefix = (
                    path_prefix.replace("\\", "\\\\")
                    .replace("%", "\\%")
                    .replace("_", "\\_")
                )
                filters.append("documents.relative_path LIKE ? ESCAPE '\\'")
                filter_parameters.append(f"{escaped_prefix}%")
            add_values_filter("documents.extension", allowed_extensions)
            add_values_filter("documents.readability", allowed_readability)
            add_values_filter("documents.indexing_state", allowed_states)
            add_values_filter("documents.source_kind", allowed_kinds)
            if modified_from:
                filters.append("documents.modified_at >= ?")
                filter_parameters.append(modified_from)
            if modified_to:
                filters.append("documents.modified_at <= ?")
                filter_parameters.append(modified_to)
            if task_document_ids is not None:
                filters.append(
                    "documents.document_id IN (SELECT value FROM json_each(?))"
                )
                filter_parameters.append(json.dumps(sorted(task_document_ids)))
            filter_sql = f"WHERE {' AND '.join(filters)}" if filters else ""
            document_count = int(
                connection.execute(
                    f"SELECT COUNT(*) FROM documents {filter_sql}",
                    filter_parameters,
                ).fetchone()[0]
            )
            metadata_sql = f"SELECT documents.* FROM documents {filter_sql}"
            metadata_parameters = list(filter_parameters)
            if document_count > 10_000:
                metadata_terms = normalized_query.split()
                metadata_clauses: list[str] = []
                metadata_values: list[object] = []
                for column in (
                    "documents.name",
                    "documents.title",
                    "documents.relative_path",
                ):
                    term_clauses = []
                    for term in metadata_terms:
                        escaped = (
                            term.replace("\\", "\\\\")
                            .replace("%", "\\%")
                            .replace("_", "\\_")
                        )
                        term_clauses.append(f"{column} LIKE ? ESCAPE '\\'")
                        metadata_values.append(f"%{escaped}%")
                    metadata_clauses.append(f"({' AND '.join(term_clauses)})")
                conjunction = " AND " if filters else " WHERE "
                metadata_sql += conjunction + f"({' OR '.join(metadata_clauses)})"
                metadata_sql += " ORDER BY relative_path COLLATE NOCASE LIMIT 2000"
                metadata_parameters.extend(metadata_values)
            rows = connection.execute(metadata_sql, metadata_parameters).fetchall()
            documents = {str(row["document_id"]): row for row in rows}
            candidates: dict[str, tuple[int, float, list[WorkspaceEvidence]]] = {}

            def allowed(row: sqlite3.Row) -> bool:
                relative_path = str(row["relative_path"])
                extension = str(row["extension"])
                if path_prefix and not relative_path.casefold().startswith(path_prefix.casefold()):
                    return False
                if allowed_extensions and extension not in allowed_extensions:
                    return False
                if (
                    allowed_readability
                    and str(row["readability"]).casefold() not in allowed_readability
                ):
                    return False
                if allowed_states and str(row["indexing_state"]).casefold() not in allowed_states:
                    return False
                row_kind = str(_row_value(row, "source_kind", "physical") or "physical")
                if allowed_kinds and row_kind.casefold() not in allowed_kinds:
                    return False
                modified_at = str(row["modified_at"])
                if modified_from and modified_at < modified_from:
                    return False
                if modified_to and modified_at > modified_to:
                    return False
                return task_document_ids is None or str(row["document_id"]) in task_document_ids

            for document_id, row in documents.items():
                relative_path = str(row["relative_path"])
                if not allowed(row):
                    continue
                filename = normalize_search_text(str(row["name"]))
                stem = normalize_search_text(Path(str(row["name"])).stem)
                title = normalize_search_text(str(row["title"]))
                path_text = normalize_search_text(relative_path)
                tier: int | None = None
                reason = ""
                if normalized_query in {filename, stem}:
                    tier, reason = 0, "文件名与查询完全一致"
                    metadata_score = 0.0
                elif normalized_query in filename:
                    tier, reason = 1, "文件名包含查询内容"
                    target = stem if normalized_query in stem else filename
                    position = target.find(normalized_query)
                    length_penalty = max(0, len(target) - len(normalized_query)) / 10_000.0
                    metadata_score = (
                        length_penalty
                        if position == 0
                        else 1.0 + (position / 1_000.0) + length_penalty
                    )
                elif (
                    normalized_query
                    and SequenceMatcher(None, normalized_query, stem).ratio() >= 0.72
                ):
                    tier, reason = 1, "文件名与查询内容相近"
                    metadata_score = 2.0 + (
                        1.0 - SequenceMatcher(None, normalized_query, stem).ratio()
                    )
                elif normalized_query in title:
                    tier, reason = 2, "文档标题包含查询内容"
                    metadata_score = 0.0
                elif normalized_query in path_text:
                    tier, reason = 2, "文件路径包含查询内容"
                    metadata_score = 0.2
                if tier is not None:
                    evidence = WorkspaceEvidence(
                        page_number=None,
                        locator=_locator_from_json(_row_value(row, "locator_json")),
                        excerpt=str(row["overview"]) or str(row["name"]),
                        reason=reason,
                        quality_score=float(row["readability_score"]),
                    )
                    candidates[document_id] = (tier, metadata_score, [evidence])
            tokens = search_terms(value)
            if tokens:
                quoted_primary = [
                    f'"{term.replace(chr(34), chr(34) * 2)}"' for term in normalized_query.split()
                ]
                fallback_expression = " OR ".join(
                    f'"{term.replace(chr(34), chr(34) * 2)}"' for term in tokens
                )
                expressions = [" AND ".join(quoted_primary)] if quoted_primary else []
                if fallback_expression not in expressions:
                    expressions.append(fallback_expression)
                matches: list[sqlite3.Row] = []
                passage_filter_sql = (
                    f"AND {' AND '.join(filters)} " if filters else ""
                )
                for expression in expressions:
                    try:
                        matches = connection.execute(
                            "SELECT passages_fts.rowid, passages_fts.document_id, "
                            "passages_fts.page_number, "
                            "passages_fts.heading, passages_fts.body, passages.locator_json, "
                            "bm25(passages_fts) AS score "
                            "FROM passages_fts JOIN passages "
                            "ON passages.passage_id = passages_fts.rowid "
                            "JOIN documents ON documents.document_id = passages_fts.document_id "
                            "WHERE passages_fts MATCH ? "
                            f"{passage_filter_sql}ORDER BY score LIMIT 250",
                            [expression, *filter_parameters],
                        ).fetchall()
                    except sqlite3.OperationalError:
                        matches = []
                    if matches:
                        break
                for match in matches:
                    document_id = str(match["document_id"])
                    if document_id not in documents:
                        lookup_filters = ["documents.document_id = ?", *filters]
                        row = connection.execute(
                            "SELECT documents.* FROM documents WHERE "
                            f"{' AND '.join(lookup_filters)}",
                            [document_id, *filter_parameters],
                        ).fetchone()
                        if row is None:
                            continue
                        documents[document_id] = row
                    row = documents[document_id]
                    relative_path = str(row["relative_path"])
                    if not allowed(row):
                        continue
                    heading = str(match["heading"])
                    body = str(match["body"])
                    query_parts = normalized_query.split()
                    normalized_heading = normalize_search_text(heading)
                    normalized_body = normalize_search_text(body)
                    heading_match = bool(query_parts) and all(
                        part in normalized_heading for part in query_parts
                    )
                    body_match = bool(query_parts) and all(
                        part in normalized_body for part in query_parts
                    )
                    if not heading_match and not body_match:
                        continue
                    tier = 2 if heading_match else 3
                    reason = "章节标题包含查询内容" if heading_match else "正文包含查询内容"
                    evidence = WorkspaceEvidence(
                        page_number=(
                            int(match["page_number"]) if match["page_number"] is not None else None
                        ),
                        heading=heading if heading_match else "",
                        locator=_locator_from_json(match["locator_json"]),
                        excerpt=_excerpt(body, value),
                        reason=reason,
                        quality_score=float(row["readability_score"]),
                    )
                    score = float(match["score"])
                    current = candidates.get(document_id)
                    if current is None:
                        candidates[document_id] = (tier, score, [evidence])
                    else:
                        best_tier, best_score, evidence_items = current
                        evidence_items.append(evidence)
                        if tier < best_tier:
                            merged_tier, merged_score = tier, score
                        elif tier == best_tier:
                            merged_tier, merged_score = best_tier, min(best_score, score)
                        else:
                            merged_tier, merged_score = best_tier, best_score
                        candidates[document_id] = (
                            merged_tier,
                            merged_score,
                            evidence_items,
                        )
            ordered = sorted(
                candidates.items(),
                key=lambda item: (
                    item[1][0],
                    item[1][1],
                    str(documents[item[0]]["name"]).casefold(),
                ),
            )[:limit]
            results: list[WorkspaceSearchResult] = []
            for rank, (document_id, (candidate_tier, _, evidence_items)) in enumerate(
                ordered, start=1
            ):
                row = documents[document_id]
                unique: list[WorkspaceEvidence] = []
                seen_evidence: set[tuple[int | None, str]] = set()
                for evidence in evidence_items:
                    key = (evidence.page_number, evidence.excerpt)
                    if key in seen_evidence:
                        continue
                    seen_evidence.add(key)
                    unique.append(evidence)
                best = unique[0]
                additional = unique[1:3]
                paged_evidence = next(
                    (evidence for evidence in unique if evidence.page_number is not None), None
                )
                if candidate_tier <= 1 and best.page_number is None and paged_evidence is not None:
                    best = paged_evidence.model_copy(
                        update={"reason": f"{best.reason}；{paged_evidence.reason}"}
                    )
                    additional = [
                        evidence
                        for evidence in unique
                        if evidence is not paged_evidence and evidence.page_number is not None
                    ][:2]
                results.append(
                    WorkspaceSearchResult(
                        document_id=document_id,
                        name=str(row["name"]),
                        relative_path=str(row["relative_path"]),
                        extension=str(row["extension"]),
                        content_hash=str(row["content_hash"]),
                        size_bytes=int(row["size_bytes"]),
                        modified_at=str(row["modified_at"]),
                        page_count=int(row["page_count"]),
                        readability=_readability_value(row["readability"]),
                        readability_score=float(row["readability_score"]),
                        indexing_state=_indexing_state_value(row["indexing_state"]),
                        source_uri=self._source_uri(row),
                        source_ref=_source_ref_from_row(row),
                        locator=best.locator,
                        quality_flags=_quality_flags(
                            _row_value(row, "quality_flags_json", "[]")
                        ),
                        error_code=str(_row_value(row, "error_code")),
                        parser_key=str(_row_value(row, "parser_key")),
                        parser_version=str(_row_value(row, "parser_version")),
                        freshness_status=_freshness_from_value(
                            _row_value(row, "freshness_status", "current")
                        ),
                        overview=str(row["overview"]),
                        best_evidence=best,
                        additional_evidence=additional,
                        rank=rank,
                    )
                )
        actual_mode: Literal["local", "assisted", "degraded"] = "local"
        degradation_reason = ""
        answer = (
            f"找到 {len(results)} 份相关资料。" if results else "当前资料空间没有找到匹配内容。"
        )
        if mode == "assisted":
            try:
                results, assisted_answer = assisted_rerank(self.workspace, value, results)
                actual_mode = "assisted"
                if assisted_answer:
                    answer = assisted_answer
            except RuntimeError as error:
                actual_mode = "degraded"
                degradation_reason = str(error) or "ai_unavailable"
            except Exception:
                actual_mode = "degraded"
                degradation_reason = "ai_unavailable"
        duration_ms = max(0, math.ceil((time.perf_counter() - started) * 1000))
        return WorkspaceSearchReport(
            query=value,
            requested_mode=mode,
            actual_mode=actual_mode,
            degradation_reason=degradation_reason,
            answer=answer,
            results=results,
            candidate_count=len(candidates),
            duration_ms=duration_ms,
        )

    def preview_path(
        self,
        document_id: str,
        page_number: int,
        highlight: str = "",
    ) -> Path:
        if page_number < 1:
            raise ValueError("Page number must be positive")
        highlight = " ".join(highlight.split()).strip()
        if len(highlight) > 200:
            raise ValueError("Preview highlight is too long")
        with closing(self._connect()) as connection:
            row = connection.execute(
                "SELECT * FROM documents WHERE document_id = ?",
                (document_id,),
            ).fetchone()
        if row is None:
            raise FileNotFoundError("Document not found")
        extension = str(row["extension"]).casefold()
        if extension in IMAGE_EXTENSIONS:
            source = self.content_path(document_id)
            destination = self.previews / str(row["content_hash"]) / "image.png"
            if not destination.exists():
                from PIL import Image

                destination.parent.mkdir(parents=True, exist_ok=True)
                with Image.open(source) as image:
                    image.convert("RGB").save(destination, format="PNG", optimize=True)
            return destination
        if extension != ".pdf":
            raise ValueError("Page preview is available only for PDF and image documents")
        if page_number > int(row["page_count"]):
            raise FileNotFoundError("Page not found")
        suffix = ""
        if highlight:
            digest = hashlib.sha256(normalize_search_text(highlight).encode("utf-8")).hexdigest()[
                :16
            ]
            suffix = f"-{digest}"
        destination = self.previews / str(row["content_hash"]) / f"{page_number}{suffix}.png"
        if destination.exists():
            return destination
        import pypdfium2 as pdfium

        source = self.content_path(document_id)
        document = pdfium.PdfDocument(str(source))
        try:
            page = document[page_number - 1]
            bitmap = page.render(scale=1.8)
            try:
                render_image: Any = bitmap.to_pil()
                if highlight:
                    from PIL import Image, ImageDraw

                    text_page = page.get_textpage()
                    boxes: list[tuple[float, float, float, float]] = []
                    try:
                        searcher = text_page.search(highlight, match_case=False)
                        try:
                            while len(boxes) < 200:
                                match = searcher.get_next()
                                if match is None:
                                    break
                                start, count = match
                                for index in range(start, start + count):
                                    try:
                                        left, bottom, right, top = text_page.get_charbox(index)
                                    except (IndexError, RuntimeError):
                                        continue
                                    if right > left and top > bottom:
                                        boxes.append((left, bottom, right, top))
                                    if len(boxes) >= 200:
                                        break
                        finally:
                            searcher.close()
                    finally:
                        text_page.close()
                    if boxes:
                        merged_boxes: list[tuple[float, float, float, float]] = []
                        for left, bottom, right, top in boxes:
                            if not merged_boxes:
                                merged_boxes.append((left, bottom, right, top))
                                continue
                            current_left, current_bottom, current_right, current_top = merged_boxes[
                                -1
                            ]
                            overlap = min(top, current_top) - max(bottom, current_bottom)
                            minimum_height = min(top - bottom, current_top - current_bottom)
                            same_line = overlap >= minimum_height * 0.45
                            close_horizontal = left - current_right <= max(6.0, minimum_height)
                            if same_line and close_horizontal:
                                merged_boxes[-1] = (
                                    min(current_left, left),
                                    min(current_bottom, bottom),
                                    max(current_right, right),
                                    max(current_top, top),
                                )
                            else:
                                merged_boxes.append((left, bottom, right, top))
                        page_width, page_height = page.get_size()
                        scale_x = render_image.width / page_width
                        scale_y = render_image.height / page_height
                        base = render_image.convert("RGBA")
                        overlay = Image.new("RGBA", base.size, (0, 0, 0, 0))
                        draw = ImageDraw.Draw(overlay)
                        for left, bottom, right, top in merged_boxes:
                            draw.rectangle(
                                (
                                    round(left * scale_x),
                                    round((page_height - top) * scale_y),
                                    round(right * scale_x),
                                    round((page_height - bottom) * scale_y),
                                ),
                                fill=(255, 211, 61, 105),
                                outline=(214, 151, 0, 190),
                                width=1,
                            )
                        render_image = Image.alpha_composite(base, overlay).convert("RGB")
                destination.parent.mkdir(parents=True, exist_ok=True)
                render_image.save(destination, format="PNG", optimize=True)
            finally:
                bitmap.close()
                page.close()
        finally:
            document.close()
        return destination

    def reprocess_document(
        self,
        document_id: str,
        progress_callback: Callable[[dict[str, Any]], None] | None = None,
    ) -> dict[str, Any]:
        with closing(self._connect()) as connection:
            row = connection.execute(
                "SELECT relative_path, source_kind, parent_document_id "
                "FROM documents WHERE document_id = ?",
                (document_id,),
            ).fetchone()
            if row is None:
                raise FileNotFoundError("Document not found")
            relative_path = str(row["relative_path"])
            force_ids = {document_id}
            if str(row["source_kind"] or "") == "archive_member":
                parent_id = str(row["parent_document_id"] or "")
                if parent_id:
                    force_ids.add(parent_id)
        result = self.sync(progress_callback, force_document_ids=force_ids)
        result["reprocessed_document_id"] = document_id
        result["relative_path"] = relative_path
        return result


def list_workspace_payloads() -> list[WorkspacePayload]:
    workspaces = ensure_v2_workspaces()
    return [WorkspaceStore(item).payload() for item in workspaces.values()]
