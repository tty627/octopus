from __future__ import annotations

from pathlib import Path

from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter

from octopus import __version__
from octopus.parsers import ParserRegistry


def test_builtin_document_parsers(tmp_path: Path) -> None:
    docx_path = tmp_path / "report.docx"
    document = Document()
    document.add_heading("Octopus Report", 1)
    document.add_paragraph("A compact project summary.")
    document.add_table(rows=2, cols=2)
    document.save(docx_path)

    xlsx_path = tmp_path / "data.xlsx"
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Summary"
    worksheet.append(["Item", "Value"])
    worksheet.append(["Octopus", "=1+1"])
    workbook.save(xlsx_path)

    pptx_path = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    textbox.text = "Octopus overview"
    presentation.save(pptx_path)

    pdf_path = tmp_path / "blank.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    with pdf_path.open("wb") as stream:
        writer.write(stream)

    image_path = tmp_path / "scan.png"
    Image.new("RGB", (16, 16), "white").save(image_path)

    registry = ParserRegistry()
    docx = registry.extract(docx_path)
    xlsx = registry.extract(xlsx_path)
    pptx = registry.extract(pptx_path)
    pdf = registry.extract(pdf_path)
    image = registry.extract(image_path)
    assert docx.document_type == "word"
    assert docx.parser_version == __version__
    assert any(item.kind == "heading" for item in docx.evidence)
    assert xlsx.metadata["sheet_names"] == ["Summary"]
    assert xlsx.extraction_stats["sampled_formulas"] == 1
    assert xlsx.evidence[0].locator == "sheet:Summary"
    assert pptx.metadata["slide_count"] == 1
    assert int(pptx.extraction_stats["shapes"]) >= 1
    assert pdf.metadata["page_count"] == 1
    assert pdf.evidence[0].locator == "page:1"
    assert image.metadata["width"] == 16
    assert image.extraction_stats["frames"] == 1


def test_unsupported_parser_emits_explicit_flag(tmp_path: Path) -> None:
    path = tmp_path / "audio.xyz"
    path.write_bytes(b"unknown")
    result = ParserRegistry().extract(path)
    assert result.unsupported
    assert result.quality_flags == ["unsupported_content_parser"]


def test_text_parser_reports_bounded_extraction(tmp_path: Path) -> None:
    path = tmp_path / "long.md"
    original = "# Heading\n" + ("content " * 12_000)
    path.write_text(original, encoding="utf-8")
    result = ParserRegistry().extract(path)
    assert result.truncated
    assert result.text_characters == len(original)
    assert result.evidence[0].kind == "heading"
    assert len(result.text) < len(original)
