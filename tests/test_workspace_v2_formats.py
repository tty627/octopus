from __future__ import annotations

import zipfile
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from octopus import workspace_v2
from octopus.workspace_tasks_v2 import WorkspaceTaskItem, create_task, load_task, save_task
from octopus.workspace_v2 import WorkspaceStore, create_workspace


def _store(raw: Path) -> WorkspaceStore:
    return WorkspaceStore(create_workspace(raw, raw.name))


def test_office_parsers_return_semantic_locators(tmp_path: Path) -> None:
    raw = tmp_path / "office"
    raw.mkdir()

    docx = Document()
    docx.add_heading("研究标题", level=1)
    docx.add_paragraph("段落中的关键证据")
    table = docx.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "方法"
    table.cell(0, 1).text = "实验"
    docx.save(raw / "paper.docx")

    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.title = "数据"
    sheet["A1"] = "指标"
    sheet["B1"] = "42"
    workbook.save(raw / "data.xlsx")

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    textbox.text = "幻灯片证据"
    presentation.save(raw / "slides.pptx")

    store = _store(raw)
    store.sync()

    results = {
        result.name: result
        for query in ("关键证据", "42", "幻灯片证据")
        for result in store.search(query).results
    }
    assert results["paper.docx"].best_evidence.locator is not None
    assert results["paper.docx"].best_evidence.locator.kind in {"paragraph", "table"}
    assert results["data.xlsx"].best_evidence.locator is not None
    assert results["data.xlsx"].best_evidence.locator.kind == "sheet"
    assert results["data.xlsx"].best_evidence.locator.sheet_name == "数据"
    assert results["slides.pptx"].best_evidence.locator is not None
    assert results["slides.pptx"].best_evidence.locator.kind == "slide"


def test_image_parser_exposes_image_locator_and_preview(tmp_path: Path, monkeypatch) -> None:
    raw = tmp_path / "images"
    raw.mkdir()
    image_path = raw / "scan.png"
    Image.new("RGB", (32, 32), "white").save(image_path)
    monkeypatch.setattr(workspace_v2, "_ocr_text", lambda image: "图片中的证据")

    store = _store(raw)
    store.sync()

    result = store.search("图片中的证据").results[0]
    assert result.best_evidence.locator is not None
    assert result.best_evidence.locator.kind == "image"
    preview = store.preview_path(result.document_id, 1)
    assert preview.exists()
    assert preview.suffix == ".png"


def test_zip_members_are_searchable_and_member_rename_requires_revalidation(tmp_path: Path) -> None:
    raw = tmp_path / "archives"
    raw.mkdir()
    archive = raw / "papers.zip"
    with zipfile.ZipFile(archive, "w", compression=zipfile.ZIP_DEFLATED) as value:
        value.writestr("paper.txt", "归档论文证据".encode())
    store = _store(raw)
    store.sync()

    result = store.search("归档论文证据").results[0]
    assert result.source_ref is not None
    assert result.source_ref.source_kind == "archive_member"
    assert result.relative_path == "papers.zip!/paper.txt"

    task = create_task(store.workspace.workspace_id, "归档核验", "验证 ZIP 成员")
    slot = task.slots[0]
    task.items.append(
        WorkspaceTaskItem(
            item_id="item-1",
            document_id=result.document_id,
            content_hash=result.content_hash,
            name=result.name,
            relative_path=result.relative_path,
            source_ref=result.source_ref,
            locator=result.best_evidence.locator,
            excerpt=result.best_evidence.excerpt,
            slot_id=slot.slot_id,
            review_state="confirmed",
        )
    )
    saved = save_task(task.workspace_id, task.task_id, task.revision, task)

    renamed = raw / "replacement.zip"
    with zipfile.ZipFile(renamed, "w", compression=zipfile.ZIP_DEFLATED) as value:
        value.writestr("renamed.txt", "归档论文证据".encode())
    archive.unlink()
    store.sync()

    from octopus.workspace_tasks_v2 import load_task

    refreshed = load_task(saved.workspace_id, saved.task_id)
    assert refreshed.items[0].source_status == "source_unconfirmed"
    assert refreshed.items[0].freshness_status == "missing"


def test_physical_duplicate_does_not_replace_archive_member_identity(tmp_path: Path) -> None:
    raw = tmp_path / "archives"
    raw.mkdir()
    payload = b"shared archived evidence " * 8
    with zipfile.ZipFile(raw / "papers.zip", "w") as archive:
        archive.writestr("paper.txt", payload)
    store = _store(raw)
    store.sync()
    member = next(
        item for item in store.list_documents() if item.source_ref.source_kind == "archive_member"
    )

    (raw / "copy.txt").write_bytes(payload)
    store.sync()

    documents = {item.relative_path: item for item in store.list_documents()}
    assert documents["papers.zip!/paper.txt"].document_id == member.document_id
    assert documents["papers.zip!/paper.txt"].source_ref.source_kind == "archive_member"
    assert documents["copy.txt"].source_ref.source_kind == "physical"
    assert documents["copy.txt"].document_id != member.document_id


def test_reprocess_archive_member_forces_container_and_member(tmp_path: Path) -> None:
    raw = tmp_path / "archives"
    raw.mkdir()
    with zipfile.ZipFile(raw / "papers.zip", "w") as archive:
        archive.writestr("paper.txt", b"archive reprocess evidence " * 8)
    store = _store(raw)
    store.sync()
    member = next(
        item for item in store.list_documents() if item.source_ref.source_kind == "archive_member"
    )

    result = store.reprocess_document(member.document_id)

    assert result["reprocessed_document_id"] == member.document_id
    assert result["indexed"] == 2
    assert result["unchanged"] == 0


def test_renaming_zip_preserves_container_member_and_task_identity(tmp_path: Path) -> None:
    raw = tmp_path / "archives"
    raw.mkdir()
    archive_path = raw / "papers.zip"
    with zipfile.ZipFile(archive_path, "w") as archive:
        archive.writestr("paper.txt", b"stable archive identity evidence " * 8)
    store = _store(raw)
    store.sync()
    documents = store.list_documents()
    container = next(item for item in documents if item.source_ref.source_kind == "archive")
    member = next(item for item in documents if item.source_ref.source_kind == "archive_member")
    task = create_task(store.workspace.workspace_id, "Archive identity")
    task.items.append(
        WorkspaceTaskItem(
            item_id="archive-item",
            document_id=member.document_id,
            content_hash=member.content_hash,
            name=member.name,
            relative_path=member.relative_path,
            source_ref=member.source_ref,
            slot_id=task.slots[0].slot_id,
            review_state="confirmed",
        )
    )
    task = save_task(task.workspace_id, task.task_id, task.revision, task)

    archive_path.rename(raw / "renamed.zip")
    store.sync()

    renamed = {item.relative_path: item for item in store.list_documents()}
    assert renamed["renamed.zip"].document_id == container.document_id
    assert renamed["renamed.zip!/paper.txt"].document_id == member.document_id
    refreshed = load_task(task.workspace_id, task.task_id)
    assert refreshed.items[0].source_status == "resolved"
    assert refreshed.items[0].freshness_status == "current"
    assert refreshed.items[0].relative_path == "renamed.zip!/paper.txt"
