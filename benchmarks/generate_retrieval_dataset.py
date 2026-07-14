from __future__ import annotations

import argparse
import json
import shutil
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from PIL import Image, ImageDraw
from pptx import Presentation

from octopus.evaluation import RETRIEVAL_DATASET_VERSION, load_retrieval_tasks


def _pdf_bytes(text: str) -> bytes:
    safe = (
        text.encode("ascii", errors="ignore")
        .decode("ascii")
        .replace("(", "[")
        .replace(")", "]")
    )
    stream = f"BT /F1 14 Tf 48 150 Td ({safe[:120]}) Tj ET".encode("ascii")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 300 200] "
        b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length "
        + str(len(stream)).encode("ascii")
        + b" >>\nstream\n"
        + stream
        + b"\nendstream",
    ]
    payload = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for number, value in enumerate(objects, start=1):
        offsets.append(len(payload))
        payload.extend(f"{number} 0 obj\n".encode("ascii"))
        payload.extend(value)
        payload.extend(b"\nendobj\n")
    xref = len(payload)
    payload.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    payload.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        payload.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    payload.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
        f"startxref\n{xref}\n%%EOF\n".encode("ascii")
    )
    return bytes(payload)


def _write_task(path: Path, kind: str, title: str, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    if kind in {"txt", "md", "py"}:
        prefix = "# " if kind == "md" else ""
        path.write_text(f"{prefix}{title}\n{content}\n", encoding="utf-8")
    elif kind == "pdf":
        path.write_bytes(_pdf_bytes(f"{title} {content}"))
    elif kind == "docx":
        document = Document()
        document.add_heading(title, level=1)
        document.add_paragraph(content)
        document.save(path)
    elif kind == "xlsx":
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Evidence"
        sheet.append(["title", "content"])
        sheet.append([title, content])
        workbook.save(path)
    elif kind == "pptx":
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
        presentation.save(path)
    elif kind == "png":
        image = Image.new("RGB", (1200, 240), "white")
        draw = ImageDraw.Draw(image)
        draw.text((30, 80), f"{title} {content}"[:140], fill="black")
        image.save(path)
    elif kind == "bin":
        path.write_bytes((f"OCTOPUS:{title}:{content}").encode())
    else:
        raise ValueError(f"Unsupported retrieval fixture format: {kind}")


def materialize_retrieval_dataset(tasks_path: Path, destination: Path) -> dict[str, object]:
    destination = destination.resolve()
    if destination.exists() and any(destination.iterdir()):
        raise ValueError(f"dataset directory must be empty: {destination}")
    destination.mkdir(parents=True, exist_ok=True)
    tasks = load_retrieval_tasks(tasks_path)
    try:
        for task in tasks:
            _write_task(
                destination.joinpath(*task.target_path.split("/")),
                task.format,
                task.title,
                task.content,
            )
    except Exception:
        shutil.rmtree(destination, ignore_errors=True)
        raise
    manifest: dict[str, object] = {
        "dataset_version": RETRIEVAL_DATASET_VERSION,
        "task_count": len(tasks),
        "files": [task.target_path for task in tasks],
    }
    (destination / "dataset-manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )
    return manifest


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("destination", type=Path)
    parser.add_argument(
        "--tasks",
        type=Path,
        default=Path(__file__).parent / "retrieval" / "v1" / "tasks.jsonl",
    )
    arguments = parser.parse_args()
    result = materialize_retrieval_dataset(arguments.tasks, arguments.destination)
    print(json.dumps(result, ensure_ascii=False))


if __name__ == "__main__":
    main()
