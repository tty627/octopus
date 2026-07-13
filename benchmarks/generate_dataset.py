from __future__ import annotations

import argparse
import hashlib
import json
import shutil
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter

KINDS = ("txt", "pdf", "docx", "xlsx", "pptx", "png", "bin")


def _templates(directory: Path) -> dict[str, Path]:
    directory.mkdir(parents=True, exist_ok=True)
    paths = {kind: directory / f"template.{kind}" for kind in KINDS if kind != "txt"}
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    with paths["pdf"].open("wb") as stream:
        writer.write(stream)
    document = Document()
    document.add_heading("Octopus benchmark", 1)
    document.add_paragraph("Deterministic non-sensitive sample document.")
    document.save(paths["docx"])
    workbook = Workbook()
    sheet = workbook.active
    sheet.append(["item", "value"])
    sheet.append(["octopus", "=1+1"])
    workbook.save(paths["xlsx"])
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    textbox.text = "Octopus benchmark"
    presentation.save(paths["pptx"])
    Image.new("RGB", (64, 32), "white").save(paths["png"])
    paths["bin"].write_bytes(bytes(range(256)))
    return paths


def generate_dataset(root: Path, count: int, mode: str = "mixed") -> dict[str, object]:
    if count < 1:
        raise ValueError("count must be positive")
    if mode not in {"mixed", "metadata"}:
        raise ValueError("mode must be mixed or metadata")
    root = root.resolve()
    if root.exists() and any(root.iterdir()):
        raise ValueError(f"dataset directory must be empty: {root}")
    root.mkdir(parents=True, exist_ok=True)
    templates = _templates(root / ".templates") if mode == "mixed" else {}
    files: list[dict[str, object]] = []
    for number in range(count):
        parts = [f"level-{depth}-{(number // (10**depth)) % 10}" for depth in range(5)]
        directory = root.joinpath(*parts)
        directory.mkdir(parents=True, exist_ok=True)
        kind = KINDS[number % len(KINDS)] if mode == "mixed" else "bin"
        name_number = number // 2 if number % 17 == 0 else number
        path = directory / f"sample-{name_number:06d}.{kind}"
        if kind == "txt":
            path.write_text(
                f"# Benchmark {number}\nOctopus deterministic evidence {number}\n",
                encoding="utf-8",
            )
        elif mode == "mixed":
            shutil.copy2(templates[kind], path)
        else:
            path.write_bytes(hashlib.sha256(str(number).encode()).digest())
        files.append(
            {
                "path": path.relative_to(root).as_posix(),
                "size": path.stat().st_size,
                "sha256": hashlib.sha256(path.read_bytes()).hexdigest(),
            }
        )
    shutil.rmtree(root / ".templates", ignore_errors=True)
    manifest: dict[str, object] = {
        "generator": "octopus-v0.3",
        "mode": mode,
        "count": count,
        "files": files,
    }
    (root / "dataset-manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    return manifest


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("root", type=Path)
    parser.add_argument("--count", type=int, default=1_000)
    parser.add_argument("--mode", choices=["mixed", "metadata"], default="mixed")
    arguments = parser.parse_args()
    result = generate_dataset(arguments.root, arguments.count, arguments.mode)
    print(json.dumps({"root": str(arguments.root.resolve()), "count": result["count"]}))


if __name__ == "__main__":
    main()
